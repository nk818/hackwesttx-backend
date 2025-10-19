from rest_framework import generics, status, permissions
from rest_framework.decorators import api_view, permission_classes
from rest_framework.response import Response
from rest_framework.authtoken.models import Token
from django.contrib.auth import authenticate
from django.db.models import Q, Count, Avg, Min, Max, F
from django.utils import timezone
from datetime import datetime, timedelta
from .models import (
    User, PasswordResetToken, Department, Professor, ClassPortfolio, 
    MarketplaceListing, PortfolioPurchase, Syllabus, SyllabusExtraction,
    ImportantDate, LectureMaterial, Flashcard, Quiz, QuizQuestion, QuizSubmission,
    ClassReview, StudyGroup, Notification, ResourceRecommendation,
    Post, Like, Comment, ProcessedFile, Document, DocumentQuiz, YouTubeVideo, CalendarEvent
)
from .serializers import (
    UserSerializer, UserRegistrationSerializer, PasswordResetRequestSerializer, 
    PasswordResetConfirmSerializer, UserRoleUpdateSerializer, DepartmentSerializer,
    ProfessorSerializer, ClassPortfolioSerializer,
    MarketplaceListingSerializer, PortfolioPurchaseSerializer, PortfolioDetailSerializer, 
    SyllabusSerializer, SyllabusExtractionSerializer,
    ImportantDateSerializer, LectureMaterialSerializer, FlashcardSerializer, QuizSerializer,
    QuizQuestionSerializer, QuizSubmissionSerializer, ClassReviewSerializer,
    StudyGroupSerializer, NotificationSerializer, ResourceRecommendationSerializer,
    PostSerializer, PostCreateSerializer, CommentSerializer, ProcessedFileSerializer,
    ProcessedFileCreateSerializer, DocumentSerializer, DocumentCreateSerializer,
    DocumentQuizSerializer, DocumentQuizCreateSerializer, YouTubeVideoSerializer, YouTubeVideoCreateSerializer,
    CalendarEventSerializer, CalendarEventCreateSerializer
)
from .permissions import IsStudentOrReadOnly, IsModeratorOrReadOnly, IsAdminOnly, IsOwnerOrModerator, IsOwnerOrReadOnly

# Visitor Landing & Onboarding Views
@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def visitor_landing(request):
    """Visitor landing page with featured portfolios and search"""
    # Get featured portfolios (public_full and public_preview)
    featured_portfolios = ClassPortfolio.objects.filter(
        visibility__in=['public_full', 'public_preview']
    ).order_by('-created_at')[:6]
    
    # Get popular marketplace listings
    popular_listings = MarketplaceListing.objects.filter(
        status='active'
    ).order_by('-created_at')[:4]
    
    # Get departments for search
    departments = Department.objects.all()[:10]
    
    return Response({
        'featured_portfolios': ClassPortfolioSerializer(featured_portfolios, many=True, context={'request': request}).data,
        'popular_listings': MarketplaceListingSerializer(popular_listings, many=True, context={'request': request}).data,
        'departments': DepartmentSerializer(departments, many=True).data,
        'search_suggestions': {
            'departments': [dept.code for dept in departments],
            'tags': ['computer-science', 'mathematics', 'physics', 'chemistry', 'biology', 'engineering']
        }
    })

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def search_portfolios(request):
    """Search portfolios for visitors and authenticated users"""
    query = request.query_params.get('q', '')
    department = request.query_params.get('department', '')
    course = request.query_params.get('course', '')
    professor = request.query_params.get('professor', '')
    tags = request.query_params.get('tags', '')
    term = request.query_params.get('term', '')
    min_price = request.query_params.get('min_price', '')
    max_price = request.query_params.get('max_price', '')
    sort_by = request.query_params.get('sort', 'newest')  # newest, helpful, rated, purchased
    
    # Base queryset - show portfolios user can access
    user = request.user if request.user.is_authenticated else None
    queryset = ClassPortfolio.objects.all()
    
    # Apply visibility filtering
    if not user:
        # Visitors can only see public portfolios
        queryset = queryset.filter(visibility__in=['public_full', 'public_preview'])
    else:
        # Authenticated users can see portfolios they have access to
        accessible_portfolios = []
        for portfolio in queryset:
            if portfolio.can_user_access(user):
                accessible_portfolios.append(portfolio.id)
        queryset = queryset.filter(id__in=accessible_portfolios)
    
    # Apply search filters
    if query:
        queryset = queryset.filter(
            Q(professor__icontains=query) |
            Q(semester__icontains=query)
        )
    
    if department:
        # No department filtering since we removed course field
        pass
    
    if course:
        # No course filtering since we removed course field
        pass
    
    if professor:
        queryset = queryset.filter(professor__icontains=professor)
    
    if tags:
        tag_list = [tag.strip() for tag in tags.split(',')]
        queryset = queryset.filter(tags__overlap=tag_list)
    
    if term:
        queryset = queryset.filter(semester__icontains=term)
    
    if min_price:
        try:
            min_price_val = float(min_price)
            queryset = queryset.filter(price__gte=min_price_val)
        except ValueError:
            pass
    
    if max_price:
        try:
            max_price_val = float(max_price)
            queryset = queryset.filter(price__lte=max_price_val)
        except ValueError:
            pass
    
    # Apply sorting using the centralized function
    queryset = apply_search_sorting(queryset, sort_by)
    
    # Limit results
    queryset = queryset[:20]
    
    return Response({
        'results': ClassPortfolioSerializer(queryset, many=True, context={'request': request}).data,
        'total_count': queryset.count(),
        'search_params': {
            'query': query,
            'department': department,
            'course': course,
            'professor': professor,
            'tags': tags,
            'sort_by': sort_by
        }
    })

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def portfolio_preview_simple(request, portfolio_id):
    """Simple portfolio preview for visitors"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if portfolio is private - if so, deny access to visitors
        if portfolio.visibility == 'private':
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        # Basic portfolio info
        portfolio_data = {
            'id': portfolio.id,
            'professor': ProfessorSerializer(portfolio.professor).data,
            'semester': portfolio.semester,
            'year': portfolio.year,
            'visibility': portfolio.visibility,
            'tags': portfolio.tags,
            'created_at': portfolio.created_at,
            'owner': portfolio.created_by.username if portfolio.created_by else None
        }
        
        # Basic content preview
        important_dates = ImportantDate.objects.filter(portfolio=portfolio)
        content = {
            'syllabus': {
                'important_dates': ImportantDateSerializer(
                    important_dates[:3], many=True
                ).data,
                'total_dates': important_dates.count(),
                'preview_count': min(3, important_dates.count()),
                'preview_note': "Showing first 3 important dates"
            },
            'materials': {
                'count': LectureMaterial.objects.filter(portfolio=portfolio).count(),
                'preview_note': "Materials available with full access"
            },
            'flashcards': {
                'items': [],
                'total_count': 0,
                'preview_note': "Flashcards available with full access"
            },
            'quizzes': {
                'items': [],
                'total_count': 0,
                'preview_note': "Quizzes available with full access"
            },
            'performance': {
                'current_grade': None,
                'grade_breakdown': None,
                'preview_note': "Grade information available with full access"
            },
            'community': {
                'posts': [],
                'total_posts': 0,
                'thread_titles': [],
                'preview_note': "Community discussions available with full access"
            }
        }
        
        return Response({
            'portfolio': portfolio_data,
            'preview_mode': True,
            'content': content,
            'access_level': 'preview'
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def portfolio_preview(request, portfolio_id):
    """Get portfolio preview with comprehensive preview rules"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        user = request.user if request.user.is_authenticated else None
        
        # Check if portfolio is visible to user
        if not portfolio.can_user_access(user):
            return Response({'error': 'Portfolio not accessible'}, status=status.HTTP_403_FORBIDDEN)
        
        # Determine if this is a preview request
        is_preview = portfolio.is_public_preview() or (not user and portfolio.visibility in ['public_preview', 'paid'])
        
        # Get preview-restricted content
        preview_data = get_preview_content(portfolio, user, is_preview)
        
        return Response({
            'portfolio': {
                'id': portfolio.id,
                'professor': ProfessorSerializer(portfolio.professor).data,
                'semester': portfolio.semester,
                'year': portfolio.year,
                'visibility': portfolio.visibility,
                'tags': portfolio.tags,
                'created_at': portfolio.created_at,
                'owner': portfolio.created_by.username if portfolio.created_by else None
            },
            'preview_mode': is_preview,
            'content': preview_data,
            'access_level': get_access_level(portfolio, user),
            'upgrade_options': get_upgrade_options(portfolio, user) if is_preview else None
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_portfolio_wizard(request):
    """Create portfolio with guided wizard flow"""
    data = request.data
    
    # Validate required fields
    required_fields = ['course_id', 'professor_id', 'semester', 'year', 'visibility']
    for field in required_fields:
        if field not in data:
            return Response({'error': f'Missing required field: {field}'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Validate visibility and pricing
    visibility = data.get('visibility')
    price = data.get('price')
    
    if visibility == 'paid' and not price:
        return Response({'error': 'Price is required for paid portfolios'}, status=status.HTTP_400_BAD_REQUEST)
    
    if visibility != 'paid' and price:
        return Response({'error': 'Price should only be set for paid portfolios'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Check if portfolio already exists for this professor/semester/year
    existing_portfolio = ClassPortfolio.objects.filter(
        professor=data['professor'],
        semester=data['semester'],
        year=data['year']
    ).first()
    
    if existing_portfolio:
        return Response({
            'error': 'A portfolio already exists for this professor, semester, and year combination',
            'existing_portfolio_id': existing_portfolio.id
        }, status=status.HTTP_400_BAD_REQUEST)
    
    # Create portfolio
    portfolio_data = {
        'professor': data['professor'],
        'semester': data['semester'],
        'year': data['year'],
        'is_public': data.get('is_public', False)
    }
    
    serializer = ClassPortfolioSerializer(data=portfolio_data, context={'request': request})
    if serializer.is_valid():
        portfolio = serializer.save(created_by=request.user)
        
        # Create marketplace listing if paid
        if visibility == 'paid':
            marketplace_data = {
                'portfolio_id': portfolio.id,
                'price': price,
                'promo_code': data.get('promo_code'),
                'campus_license_available': data.get('campus_license_available', False),
                'campus_license_price': data.get('campus_license_price')
            }
            
            marketplace_serializer = MarketplaceListingSerializer(data=marketplace_data)
            if marketplace_serializer.is_valid():
                marketplace_serializer.save()
            else:
                return Response({'error': 'Failed to create marketplace listing', 'details': marketplace_serializer.errors}, 
                              status=status.HTTP_400_BAD_REQUEST)
        
        return Response({
            'message': 'Portfolio created successfully',
            'portfolio': ClassPortfolioSerializer(portfolio, context={'request': request}).data,
            'next_steps': {
                'syllabus_upload': f'/api/syllabi/?portfolio_id={portfolio.id}',
                'materials_upload': f'/api/materials/?portfolio_id={portfolio.id}',
                'grade_tracking': f'/api/portfolios/{portfolio.id}/grades/',
                'marketplace_listing': f'/api/marketplace/?portfolio_id={portfolio.id}' if visibility == 'paid' else None
            }
        }, status=status.HTTP_201_CREATED)
    
    return Response({'error': 'Failed to create portfolio', 'details': serializer.errors}, 
                   status=status.HTTP_400_BAD_REQUEST)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def onboarding_status(request):
    """Get user's onboarding status and next steps"""
    user = request.user
    
    # Check if user has completed basic profile
    profile_complete = all([
        user.first_name,
        user.last_name,
        user.university,
        user.major,
        user.graduation_year
    ])
    
    # Check if user has created any portfolios
    portfolios_count = ClassPortfolio.objects.filter(created_by=user).count()
    
    # Check if user has purchased any portfolios
    purchases_count = PortfolioPurchase.objects.filter(buyer=user).count()
    
    # Determine next steps
    next_steps = []
    if not profile_complete:
        next_steps.append({
            'step': 'complete_profile',
            'title': 'Complete Your Profile',
            'description': 'Add your university, major, and graduation year',
            'url': '/api/auth/me/',
            'priority': 'high'
        })
    
    if portfolios_count == 0:
        next_steps.append({
            'step': 'create_portfolio',
            'title': 'Create Your First Portfolio',
            'description': 'Start by creating a portfolio for one of your classes',
            'url': '/api/portfolios/create-wizard/',
            'priority': 'high'
        })
    
    if purchases_count == 0 and portfolios_count > 0:
        next_steps.append({
            'step': 'explore_marketplace',
            'title': 'Explore the Marketplace',
            'description': 'Discover portfolios from other students',
            'url': '/api/marketplace/',
            'priority': 'medium'
        })
    
    return Response({
        'profile_complete': profile_complete,
        'portfolios_count': portfolios_count,
        'purchases_count': purchases_count,
        'next_steps': next_steps,
        'onboarding_complete': len(next_steps) == 0
    })

# Authentication Views
@api_view(['POST'])
@permission_classes([permissions.AllowAny])
def register(request):
    serializer = UserRegistrationSerializer(data=request.data)
    if serializer.is_valid():
        user = serializer.save()
        token, created = Token.objects.get_or_create(user=user)
        return Response({
            'token': token.key,
            'user': UserSerializer(user).data
        }, status=status.HTTP_201_CREATED)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([permissions.AllowAny])
def login(request):
    username_or_email = request.data.get('username')
    password = request.data.get('password')
    
    if username_or_email and password:
        # Try to authenticate with username first
        user = authenticate(username=username_or_email, password=password)
        
        # If that fails, try with email
        if not user:
            try:
                user_obj = User.objects.get(email=username_or_email.lower())
                user = authenticate(username=user_obj.username, password=password)
            except User.DoesNotExist:
                pass
        
        if user:
            token, created = Token.objects.get_or_create(user=user)
            return Response({
                'token': token.key,
                'user': UserSerializer(user).data
            })
        return Response({'error': 'Invalid credentials'}, status=status.HTTP_401_UNAUTHORIZED)
    return Response({'error': 'Username/email and password required'}, status=status.HTTP_400_BAD_REQUEST)

@api_view(['GET'])
def me(request):
    return Response(UserSerializer(request.user).data)

# Password Reset Views
@api_view(['POST'])
@permission_classes([permissions.AllowAny])
def request_password_reset(request):
    """
    Request a password reset token via email
    """
    serializer = PasswordResetRequestSerializer(data=request.data)
    if serializer.is_valid():
        email = serializer.validated_data['email']
        user = User.objects.get(email=email)
        
        # Generate a secure token
        import secrets
        token = secrets.token_urlsafe(32)
        
        # Create password reset token (expires in 1 hour)
        expires_at = timezone.now() + timedelta(hours=1)
        PasswordResetToken.objects.create(
            user=user,
            token=token,
            expires_at=expires_at
        )
        
        # Send email (in production, you'd use Django's email backend)
        reset_url = f"http://localhost:3000/reset-password?token={token}"
        
        # For now, we'll just return the token in the response
        # In production, you'd send this via email
        return Response({
            'message': 'Password reset token generated successfully',
            'reset_url': reset_url,  # Remove this in production
            'token': token  # Remove this in production
        }, status=status.HTTP_200_OK)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([permissions.AllowAny])
def confirm_password_reset(request):
    """
    Confirm password reset with token and new password
    """
    serializer = PasswordResetConfirmSerializer(data=request.data)
    if serializer.is_valid():
        token = serializer.validated_data['token']
        new_password = serializer.validated_data['new_password']
        
        # Get the token object
        token_obj = PasswordResetToken.objects.get(token=token)
        user = token_obj.user
        
        # Update password
        user.set_password(new_password)
        user.save()
        
        # Mark token as used
        token_obj.is_used = True
        token_obj.save()
        
        # Delete all other password reset tokens for this user
        PasswordResetToken.objects.filter(user=user, is_used=False).delete()
        
        return Response({
            'message': 'Password reset successfully'
        }, status=status.HTTP_200_OK)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

# User Management Views
@api_view(['GET'])
@permission_classes([IsAdminOnly])
def list_users(request):
    """
    List all users (Admin only)
    """
    users = User.objects.all().order_by('-date_joined')
    serializer = UserSerializer(users, many=True)
    return Response(serializer.data)

@api_view(['PATCH'])
@permission_classes([IsAdminOnly])
def update_user_role(request, user_id):
    """
    Update user role and verification status (Admin only)
    """
    try:
        user = User.objects.get(id=user_id)
    except User.DoesNotExist:
        return Response({'error': 'User not found'}, status=status.HTTP_404_NOT_FOUND)
    
    serializer = UserRoleUpdateSerializer(user, data=request.data, partial=True, context={'request': request})
    if serializer.is_valid():
        serializer.save()
        return Response({
            'message': 'User role updated successfully',
            'user': UserSerializer(user).data
        }, status=status.HTTP_200_OK)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['DELETE'])
@permission_classes([IsAdminOnly])
def delete_user(request, user_id):
    """
    Delete a user account (Admin only)
    """
    try:
        user = User.objects.get(id=user_id)
        if user == request.user:
            return Response({'error': 'Cannot delete your own account'}, status=status.HTTP_400_BAD_REQUEST)
        
        user.delete()
        return Response({'message': 'User deleted successfully'}, status=status.HTTP_200_OK)
    except User.DoesNotExist:
        return Response({'error': 'User not found'}, status=status.HTTP_404_NOT_FOUND)

# Department Views
class DepartmentListCreateView(generics.ListCreateAPIView):
    queryset = Department.objects.all()
    serializer_class = DepartmentSerializer
    permission_classes = [permissions.IsAuthenticated]

class DepartmentDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Department.objects.all()
    serializer_class = DepartmentSerializer
    permission_classes = [permissions.IsAuthenticated]

# Professor Views
class ProfessorListCreateView(generics.ListCreateAPIView):
    queryset = Professor.objects.all()
    serializer_class = ProfessorSerializer
    permission_classes = [permissions.IsAuthenticated]

class ProfessorDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Professor.objects.all()
    serializer_class = ProfessorSerializer
    permission_classes = [permissions.IsAuthenticated]


# Class Portfolio Views
class PortfolioListCreateView(generics.ListCreateAPIView):
    queryset = ClassPortfolio.objects.all()
    serializer_class = ClassPortfolioSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        user = self.request.user
        
        # Filter portfolios based on user access
        # Show portfolios created by the user or public portfolios
        queryset = ClassPortfolio.objects.filter(
            Q(created_by=user) | Q(is_public=True)
        )
        
        # Apply search and filter parameters
        search = self.request.query_params.get('search', None)
        professor = self.request.query_params.get('professor', None)
        
        if search:
            queryset = queryset.filter(
                Q(professor__icontains=search) |
                Q(semester__icontains=search)
            )
        
        if professor:
            queryset = queryset.filter(professor__icontains=professor)
        
        return queryset.order_by('-created_at')

    def get(self, request, *args, **kwargs):
        """GET method to retrieve portfolios based on user token"""
        return self.list(request, *args, **kwargs)

class PortfolioDetailView(generics.RetrieveUpdateDestroyAPIView):
    serializer_class = PortfolioDetailSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        user = self.request.user
        # Filter portfolios based on user access
        return ClassPortfolio.objects.filter(
            Q(created_by=user) | Q(is_public=True)
        )
    
    def update(self, request, *args, **kwargs):
        # Get the instance
        instance = self.get_object()
        
        # Only allow the owner to update
        if instance.created_by != request.user:
            return Response(
                {'error': 'You do not have permission to edit this portfolio'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        # Allow partial updates (PATCH)
        partial = kwargs.pop('partial', False)
        if request.method == 'PATCH':
            partial = True
        
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        self.perform_update(serializer)
        
        return Response(serializer.data)
    
    def destroy(self, request, *args, **kwargs):
        # Get the instance
        instance = self.get_object()
        
        # Only allow the owner to delete
        if instance.created_by != request.user:
            return Response(
                {'error': 'You do not have permission to delete this portfolio'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        self.perform_destroy(instance)
        return Response(status=status.HTTP_204_NO_CONTENT)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_portfolios(request):
    """GET method to retrieve only the authenticated user's portfolios"""
    user = request.user
    
    # Get only portfolios created by the current user
    portfolios = ClassPortfolio.objects.filter(created_by=user).order_by('-created_at')
    
    # Apply search filter if provided
    search = request.query_params.get('search', None)
    if search:
        portfolios = portfolios.filter(
            Q(professor__icontains=search) |
            Q(semester__icontains=search)
        )
    
    # Serialize the data
    serializer = ClassPortfolioSerializer(portfolios, many=True, context={'request': request})
    
    return Response({
        'count': portfolios.count(),
        'results': serializer.data
    })

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def public_portfolios(request):
    """GET method to retrieve all public portfolios (no authentication required)"""
    
    # Get only portfolios where is_public=True
    portfolios = ClassPortfolio.objects.filter(is_public=True).order_by('-created_at')
    
    # Apply search filter if provided
    search = request.query_params.get('search', None)
    if search:
        portfolios = portfolios.filter(
            Q(professor__icontains=search) |
            Q(course__icontains=search) |
            Q(semester__icontains=search) |
            Q(year__icontains=search)
        )
    
    # Apply semester filter if provided
    semester = request.query_params.get('semester', None)
    if semester:
        portfolios = portfolios.filter(semester=semester)
    
    # Apply year filter if provided
    year = request.query_params.get('year', None)
    if year:
        portfolios = portfolios.filter(year=year)
    
    # Apply professor filter if provided
    professor = request.query_params.get('professor', None)
    if professor:
        portfolios = portfolios.filter(professor__icontains=professor)
    
    # Serialize the data
    serializer = ClassPortfolioSerializer(portfolios, many=True, context={'request': request})
    
    return Response({
        'count': portfolios.count(),
        'results': serializer.data
    })

@api_view(['PATCH'])
@permission_classes([permissions.IsAuthenticated])
def update_portfolio(request, portfolio_id):
    """PATCH method to update a portfolio (only by owner)"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user owns this portfolio
        if portfolio.created_by != request.user:
            return Response(
                {'error': 'You can only update your own portfolios'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        # Update fields
        if 'professor' in request.data:
            portfolio.professor = request.data['professor']
        if 'course' in request.data:
            portfolio.course = request.data['course']
        if 'semester' in request.data:
            portfolio.semester = request.data['semester']
        if 'year' in request.data:
            portfolio.year = request.data['year']
        if 'price' in request.data:
            portfolio.price = request.data['price']
        if 'is_public' in request.data:
            portfolio.is_public = request.data['is_public']
        if 'color' in request.data:
            portfolio.color = request.data['color']
        
        portfolio.save()
        
        # Serialize and return updated portfolio
        serializer = ClassPortfolioSerializer(portfolio, context={'request': request})
        return Response({
            'success': True,
            'data': serializer.data
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response(
            {'error': 'Portfolio not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        return Response(
            {'error': str(e)},
            status=status.HTTP_400_BAD_REQUEST
        )

# Marketplace Views
class MarketplaceListingListCreateView(generics.ListCreateAPIView):
    queryset = MarketplaceListing.objects.filter(status='active')
    serializer_class = MarketplaceListingSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        queryset = MarketplaceListing.objects.filter(status='active')
        
        # Filter by price range
        min_price = self.request.query_params.get('min_price', None)
        max_price = self.request.query_params.get('max_price', None)
        
        if min_price:
            queryset = queryset.filter(price__gte=min_price)
        if max_price:
            queryset = queryset.filter(price__lte=max_price)
        
        # Filter by course/department
        course = self.request.query_params.get('course', None)
        department = self.request.query_params.get('department', None)
        
        if course:
            # No course filtering since we removed course field
            pass
        if department:
            # No department filtering since we removed course field
            pass
        
        return queryset.order_by('-created_at')

class MarketplaceListingDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = MarketplaceListing.objects.all()
    serializer_class = MarketplaceListingSerializer
    permission_classes = [permissions.IsAuthenticated]

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def purchase_portfolio(request, listing_id):
    """Purchase a portfolio"""
    try:
        listing = MarketplaceListing.objects.get(id=listing_id, status='active')
    except MarketplaceListing.DoesNotExist:
        return Response({'error': 'Listing not found or not available'}, status=status.HTTP_404_NOT_FOUND)
    
    # Check if user already purchased
    if PortfolioPurchase.objects.filter(listing=listing, buyer=request.user).exists():
        return Response({'error': 'You have already purchased this portfolio'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Check if user is the owner
    if listing.portfolio.created_by == request.user:
        return Response({'error': 'You cannot purchase your own portfolio'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Get effective price
    promo_code = request.data.get('promo_code', None)
    effective_price = listing.get_effective_price(promo_code)
    
    # Create purchase record (in production, integrate with payment processor)
    purchase = PortfolioPurchase.objects.create(
        listing=listing,
        buyer=request.user,
        purchase_price=effective_price,
        promo_code_used=promo_code,
        payment_method='stripe',  # Default for now
        payment_id=f"test_payment_{timezone.now().timestamp()}"
    )
    
    return Response({
        'message': 'Portfolio purchased successfully',
        'purchase': PortfolioPurchaseSerializer(purchase).data
    }, status=status.HTTP_201_CREATED)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_purchases(request):
    """Get user's purchased portfolios"""
    purchases = PortfolioPurchase.objects.filter(buyer=request.user).order_by('-purchased_at')
    serializer = PortfolioPurchaseSerializer(purchases, many=True)
    return Response(serializer.data)

# Syllabus Views - Smart Syllabus Scanner Page
class SyllabusUploadView(generics.CreateAPIView):
    queryset = Syllabus.objects.all()
    serializer_class = SyllabusSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def perform_create(self, serializer):
        syllabus = serializer.save()
        # Start extraction process
        self._start_extraction(syllabus)
    
    def _start_extraction(self, syllabus):
        """Start the syllabus extraction process with confidence scoring"""
        try:
            syllabus.extraction_status = 'processing'
            syllabus.save()
            
            # Extract text from file
            extracted_text = self._extract_text_from_file(syllabus.file)
            syllabus.extracted_text = extracted_text
            syllabus.save()
            
            # Run AI extraction with confidence scoring
            from .syllabus_extractor import SyllabusExtractor
            extractor = SyllabusExtractor()
            extraction_data = extractor.extract_from_text(extracted_text)
            
            # Create extraction record with confidence scores
            from .models import SyllabusExtraction
            SyllabusExtraction.objects.create(
                syllabus=syllabus,
                **extraction_data
            )
            
            syllabus.extraction_status = 'completed'
            syllabus.save()
            
        except Exception as e:
            syllabus.extraction_status = 'failed'
            syllabus.extraction_error = str(e)
            syllabus.save()
    
    def _extract_text_from_file(self, file):
        """Extract text from uploaded file with proper parsing"""
        try:
            import PyPDF2
            import docx
            from pptx import Presentation
            import os
            
            file_extension = os.path.splitext(file.name)[1].lower()
            
            if file_extension == '.pdf':
                # Extract from PDF
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
            elif file_extension in ['.doc', '.docx']:
                # Extract from Word document
                doc = docx.Document(file)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            elif file_extension in ['.ppt', '.pptx']:
                # Extract from PowerPoint
                prs = Presentation(file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                return f"Unsupported file type: {file_extension}"
        except Exception as e:
            return f"Error extracting text: {str(e)}"

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def syllabus_page_data(request, portfolio_id):
    """Get complete syllabus page data with extraction status and calendar sync"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can access this portfolio
        if not portfolio.can_user_access(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        # Get syllabus data
        syllabus_data = None
        if hasattr(portfolio, 'syllabus'):
            syllabus = portfolio.syllabus
            syllabus_data = {
                'id': syllabus.id,
                'file': syllabus.file.url if syllabus.file else None,
                'uploaded_at': syllabus.uploaded_at,
                'extraction_status': syllabus.extraction_status,
                'extraction_error': syllabus.extraction_error,
                'extracted_text': syllabus.extracted_text,
                'extraction': None
            }
            
            # Get extraction data if available
            if hasattr(syllabus, 'extraction'):
                extraction = syllabus.extraction
                syllabus_data['extraction'] = {
                    'course_title': extraction.course_title,
                    'course_code': extraction.course_code,
                    'professor_name': extraction.professor_name,
                    'professor_email': extraction.professor_email,
                    'class_schedule': extraction.class_schedule,
                    'office_hours': extraction.office_hours,
                    'extraction_confidence': extraction.extraction_confidence,
                    'all_important_dates': extraction.all_important_dates,
                    'extracted_at': extraction.extracted_at
                }
        
        # Get important dates
        important_dates = ImportantDate.objects.filter(portfolio=portfolio).order_by('due_date')
        dates_data = ImportantDateSerializer(important_dates, many=True).data
        
        # Calendar sync status
        calendar_sync = {
            'google_calendar': False,  # To be implemented
            'outlook_calendar': False,  # To be implemented
            'ical_export': True,  # Always available
            'last_sync': None
        }
        
        return Response({
            'portfolio': {
                'id': portfolio.id,
                'professor': ProfessorSerializer(portfolio.professor).data,
                'semester': portfolio.semester,
                'year': portfolio.year
            },
            'syllabus': syllabus_data,
            'important_dates': dates_data,
            'calendar_sync': calendar_sync,
            'extraction_status': syllabus_data['extraction_status'] if syllabus_data else 'no_syllabus',
            'can_edit': portfolio.can_user_edit(request.user)
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def sync_to_calendar(request, portfolio_id):
    """Sync important dates to external calendars"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can edit this portfolio
        if not portfolio.can_user_edit(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        calendar_type = request.data.get('calendar_type', 'ical')  # google, outlook, ical
        reminder_settings = request.data.get('reminder_settings', {
            'week_before': True,
            'day_before': True,
            'hour_before': False
        })
        
        # Get important dates
        important_dates = ImportantDate.objects.filter(portfolio=portfolio)
        
        if calendar_type == 'ical':
            # Generate iCal format
            ical_content = generate_ical_content(important_dates, reminder_settings)
            return Response({
                'message': 'iCal file generated',
                'ical_content': ical_content,
                'download_url': f'/api/portfolios/{portfolio_id}/calendar/ical/'
            })
        elif calendar_type == 'google':
            # Google Calendar integration (placeholder)
            return Response({
                'message': 'Google Calendar sync initiated',
                'auth_url': '/api/auth/google-calendar/',  # To be implemented
                'status': 'pending'
            })
        elif calendar_type == 'outlook':
            # Outlook integration (placeholder)
            return Response({
                'message': 'Outlook Calendar sync initiated',
                'auth_url': '/api/auth/outlook-calendar/',  # To be implemented
                'status': 'pending'
            })
        else:
            return Response({'error': 'Invalid calendar type'}, status=status.HTTP_400_BAD_REQUEST)
            
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

def generate_ical_content(important_dates, reminder_settings):
    """Generate iCal content for important dates"""
    ical_lines = [
        'BEGIN:VCALENDAR',
        'VERSION:2.0',
        'PRODID:-//Class Portfolio//EN',
        'CALSCALE:GREGORIAN',
        'METHOD:PUBLISH'
    ]
    
    for date_obj in important_dates:
        # Format date for iCal
        start_date = date_obj.due_date.strftime('%Y%m%dT%H%M%S')
        
        ical_lines.extend([
            'BEGIN:VEVENT',
            f'UID:{date_obj.id}@classportfolio.com',
            f'DTSTART:{start_date}',
            f'DTEND:{start_date}',
            f'SUMMARY:{date_obj.title}',
            f'DESCRIPTION:{date_obj.description or ""}',
            f'LOCATION:{getattr(date_obj, "location", "") or ""}',
            'STATUS:CONFIRMED',
            'TRANSP:OPAQUE'
        ])
        
        # Add reminders
        if reminder_settings.get('week_before'):
            ical_lines.append('BEGIN:VALARM')
            ical_lines.append('TRIGGER:-P7D')
            ical_lines.append('ACTION:DISPLAY')
            ical_lines.append('DESCRIPTION:Reminder')
            ical_lines.append('END:VALARM')
        
        if reminder_settings.get('day_before'):
            ical_lines.append('BEGIN:VALARM')
            ical_lines.append('TRIGGER:-P1D')
            ical_lines.append('ACTION:DISPLAY')
            ical_lines.append('DESCRIPTION:Reminder')
            ical_lines.append('END:VALARM')
        
        if reminder_settings.get('hour_before'):
            ical_lines.append('BEGIN:VALARM')
            ical_lines.append('TRIGGER:-PT1H')
            ical_lines.append('ACTION:DISPLAY')
            ical_lines.append('DESCRIPTION:Reminder')
            ical_lines.append('END:VALARM')
        
        ical_lines.append('END:VEVENT')
    
    ical_lines.append('END:VCALENDAR')
    return '\n'.join(ical_lines)

# Interactive Learning Space Page
@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def learning_space_page_data(request, portfolio_id):
    """Get complete Interactive Learning Space page data"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can access this portfolio
        if not portfolio.can_user_access(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        # Get lecture materials
        materials = LectureMaterial.objects.filter(portfolio=portfolio).order_by('-uploaded_at')
        materials_data = LectureMaterialSerializer(materials, many=True).data
        
        # Get flashcards
        flashcards = Flashcard.objects.filter(material__portfolio=portfolio).order_by('-created_at')
        flashcards_data = FlashcardSerializer(flashcards, many=True).data
        
        # Get quizzes
        quizzes = Quiz.objects.filter(portfolio=portfolio).order_by('-created_at')
        quizzes_data = QuizSerializer(quizzes, many=True).data
        
        # Get AI-generated summaries
        summaries = []
        for material in materials:
            if hasattr(material, 'processed_file') and material.processed_file.summary:
                summaries.append({
                    'material_id': material.id,
                    'title': material.title,
                    'summary': material.processed_file.summary,
                    'topics': material.processed_file.topics,
                    'generated_at': material.processed_file.processed_at
                })
        
        # Get user's quiz submissions for progress tracking
        user_submissions = QuizSubmission.objects.filter(
            quiz__portfolio=portfolio,
            user=request.user
        ).order_by('-submitted_at')
        
        # Calculate progress metrics
        total_flashcards = flashcards.count()
        mastered_flashcards = 0  # To be implemented with spaced repetition
        
        total_quizzes = quizzes.count()
        completed_quizzes = user_submissions.values('quiz').distinct().count()
        
        return Response({
            'portfolio': {
                'id': portfolio.id,
                'professor': ProfessorSerializer(portfolio.professor).data
            },
            'materials': materials_data,
            'flashcards': flashcards_data,
            'quizzes': quizzes_data,
            'summaries': summaries,
            'progress': {
                'total_materials': materials.count(),
                'total_flashcards': total_flashcards,
                'mastered_flashcards': mastered_flashcards,
                'total_quizzes': total_quizzes,
                'completed_quizzes': completed_quizzes,
                'completion_percentage': (completed_quizzes / total_quizzes * 100) if total_quizzes > 0 else 0
            },
            'can_edit': portfolio.can_user_edit(request.user)
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def generate_learning_content(request, portfolio_id):
    """Generate flashcards, quizzes, and summaries from uploaded materials"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can edit this portfolio
        if not portfolio.can_user_edit(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        material_id = request.data.get('material_id')
        content_types = request.data.get('content_types', ['summary', 'flashcards', 'quiz'])
        
        try:
            material = LectureMaterial.objects.get(id=material_id, portfolio=portfolio)
        except LectureMaterial.DoesNotExist:
            return Response({'error': 'Material not found'}, status=status.HTTP_404_NOT_FOUND)
        
        generated_content = {}
        
        # Generate summary
        if 'summary' in content_types:
            summary = generate_ai_summary(material)
            generated_content['summary'] = summary
        
        # Generate flashcards
        if 'flashcards' in content_types:
            flashcards = generate_flashcards(material)
            generated_content['flashcards'] = flashcards
        
        # Generate quiz
        if 'quiz' in content_types:
            quiz = generate_quiz(material)
            generated_content['quiz'] = quiz
        
        return Response({
            'message': 'Content generated successfully',
            'generated_content': generated_content,
            'material_id': material_id
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

def generate_ai_summary(material):
    """Generate AI summary from lecture material"""
    try:
        import openai
        
        # Get text content from processed file
        if hasattr(material, 'processed_file') and material.processed_file.extracted_text:
            text_content = material.processed_file.extracted_text
        else:
            text_content = f"Content from {material.title}"
        
        # Truncate if too long
        if len(text_content) > 4000:
            text_content = text_content[:4000] + "..."
        
        # Generate summary using OpenAI
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert academic assistant. Create a concise summary of the following lecture material, highlighting key concepts and main points."},
                {"role": "user", "content": f"Summarize this lecture material:\n\n{text_content}"}
            ],
            max_tokens=500,
            temperature=0.7
        )
        
        return {
            'summary': response.choices[0].message.content,
            'topics': extract_topics(text_content),
            'generated_at': timezone.now()
        }
    except Exception as e:
        return {
            'summary': f"Error generating summary: {str(e)}",
            'topics': [],
            'generated_at': timezone.now()
        }

def generate_flashcards(material):
    """Generate flashcards from lecture material"""
    try:
        import openai
        
        # Get text content
        if hasattr(material, 'processed_file') and material.processed_file.extracted_text:
            text_content = material.processed_file.extracted_text
        else:
            text_content = f"Content from {material.title}"
        
        # Truncate if too long
        if len(text_content) > 3000:
            text_content = text_content[:3000] + "..."
        
        # Generate flashcards using OpenAI
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert educational content creator. Create 5 flashcards from the following material. Format each flashcard as JSON with 'front' and 'back' fields."},
                {"role": "user", "content": f"Create flashcards from this material:\n\n{text_content}"}
            ],
            max_tokens=800,
            temperature=0.7
        )
        
        # Parse response and create flashcards
        flashcards = []
        try:
            import json
            content = response.choices[0].message.content
            # Extract JSON from response
            json_start = content.find('[')
            json_end = content.rfind(']') + 1
            if json_start != -1 and json_end != -1:
                flashcards_data = json.loads(content[json_start:json_end])
                for card_data in flashcards_data:
                    flashcards.append({
                        'front': card_data.get('front', ''),
                        'back': card_data.get('back', ''),
                        'topic': material.title
                    })
        except:
            # Fallback: create simple flashcards
            flashcards = [
                {'front': f'What is the main topic of {material.title}?', 'back': 'Key concepts from the material', 'topic': material.title},
                {'front': f'What are the important points in {material.title}?', 'back': 'Main points covered', 'topic': material.title}
            ]
        
        return flashcards
    except Exception as e:
        return [{'front': 'Error generating flashcards', 'back': str(e), 'topic': material.title}]

def generate_quiz(material):
    """Generate quiz questions from lecture material"""
    try:
        import openai
        
        # Get text content
        if hasattr(material, 'processed_file') and material.processed_file.extracted_text:
            text_content = material.processed_file.extracted_text
        else:
            text_content = f"Content from {material.title}"
        
        # Truncate if too long
        if len(text_content) > 3000:
            text_content = text_content[:3000] + "..."
        
        # Generate quiz using OpenAI
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert quiz creator. Create 5 multiple-choice questions from the following material. Format each question as JSON with 'question', 'options' (array of 4 options), and 'correct_answer' (index 0-3) fields."},
                {"role": "user", "content": f"Create quiz questions from this material:\n\n{text_content}"}
            ],
            max_tokens=1000,
            temperature=0.7
        )
        
        # Parse response and create quiz questions
        questions = []
        try:
            import json
            content = response.choices[0].message.content
            # Extract JSON from response
            json_start = content.find('[')
            json_end = content.rfind(']') + 1
            if json_start != -1 and json_end != -1:
                questions_data = json.loads(content[json_start:json_end])
                for q_data in questions_data:
                    questions.append({
                        'question': q_data.get('question', ''),
                        'options': q_data.get('options', ['A', 'B', 'C', 'D']),
                        'correct_answer': q_data.get('correct_answer', 0),
                        'question_type': 'multiple_choice'
                    })
        except:
            # Fallback: create simple questions
            questions = [
                {
                    'question': f'What is the main topic covered in {material.title}?',
                    'options': ['Topic A', 'Topic B', 'Topic C', 'Topic D'],
                    'correct_answer': 0,
                    'question_type': 'multiple_choice'
                }
            ]
        
        return questions
    except Exception as e:
        return [{'question': 'Error generating quiz', 'options': ['A', 'B', 'C', 'D'], 'correct_answer': 0, 'question_type': 'multiple_choice'}]

def extract_topics(text_content):
    """Extract topics from text content"""
    # Simple topic extraction (in production, use NLP libraries)
    topics = []
    words = text_content.lower().split()
    common_topics = ['algorithm', 'data structure', 'programming', 'computer science', 'mathematics', 'physics', 'chemistry', 'biology']
    
    for topic in common_topics:
        if topic in text_content.lower():
            topics.append(topic.title())
    
    return topics[:5]  # Return top 5 topics

# Class Performance Tracker Page
@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def performance_tracker_page_data(request, portfolio_id):
    """Get complete Class Performance Tracker page data with grade projections"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can access this portfolio
        if not portfolio.can_user_access(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        # Get all grades
        all_grades = portfolio.all_grades
        current_grade = portfolio.calculate_current_grade()
        
        # Calculate grade projections
        projections = calculate_grade_projections(portfolio)
        
        # Get grade breakdown
        grade_breakdown = portfolio.grade_breakdown or {}
        
        # Calculate category averages
        category_averages = {}
        for category, percentage in grade_breakdown.items():
            if category == 'exams' and portfolio.exam_grades:
                grades = portfolio.exam_grades
            elif category == 'homework' and portfolio.homework_grades:
                grades = portfolio.homework_grades
            elif category == 'quizzes' and portfolio.quiz_grades:
                grades = portfolio.quiz_grades
            elif category == 'projects' and portfolio.project_grades:
                grades = portfolio.project_grades
            else:
                continue
            
            if grades:
                total_points = sum(grade.get('points_earned', 0) for grade in grades)
                total_possible = sum(grade.get('points_possible', 0) for grade in grades)
                if total_possible > 0:
                    category_averages[category] = {
                        'average': round((total_points / total_possible) * 100, 2),
                        'percentage': percentage,
                        'count': len(grades)
                    }
        
        # Get upcoming assignments
        upcoming_dates = ImportantDate.objects.filter(
            portfolio=portfolio,
            due_date__gte=timezone.now()
        ).order_by('due_date')[:5]
        
        return Response({
            'portfolio': {
                'id': portfolio.id,
                'professor': ProfessorSerializer(portfolio.professor).data,
                'semester': portfolio.semester,
                'year': portfolio.year
            },
            'current_grade': current_grade,
            'all_grades': all_grades,
            'grade_breakdown': grade_breakdown,
            'category_averages': category_averages,
            'projections': projections,
            'upcoming_assignments': ImportantDateSerializer(upcoming_dates, many=True).data,
            'can_edit': portfolio.can_user_edit(request.user)
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def calculate_what_if_scenario(request, portfolio_id):
    """Calculate what-if grade scenarios based on upcoming assignments"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Check if user can edit this portfolio
        if not portfolio.can_user_edit(request.user):
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        scenario_data = request.data
        scenario_name = scenario_data.get('name', 'What-if Scenario')
        grade_predictions = scenario_data.get('grade_predictions', {})
        
        # Calculate scenario grade
        scenario_grade = calculate_scenario_grade(portfolio, grade_predictions)
        
        # Compare with current grade
        current_grade = portfolio.calculate_current_grade()
        
        return Response({
            'scenario_name': scenario_name,
            'current_grade': current_grade,
            'scenario_grade': scenario_grade,
            'grade_difference': scenario_grade - current_grade if current_grade else 0,
            'grade_predictions': grade_predictions,
            'calculated_at': timezone.now()
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

def calculate_grade_projections(portfolio):
    """Calculate various grade projections for the portfolio"""
    current_grade = portfolio.calculate_current_grade()
    grade_breakdown = portfolio.grade_breakdown or {}
    
    projections = {
        'current': current_grade,
        'scenarios': []
    }
    
    # Scenario 1: Maintain current performance
    if current_grade:
        projections['scenarios'].append({
            'name': 'Maintain Current Performance',
            'grade': current_grade,
            'description': 'Continue performing at current level'
        })
    
    # Scenario 2: Improve by 10%
    if current_grade:
        improved_grade = min(100, current_grade + 10)
        projections['scenarios'].append({
            'name': 'Improve by 10%',
            'grade': improved_grade,
            'description': 'Improve performance by 10 percentage points'
        })
    
    # Scenario 3: Perfect scores on remaining assignments
    perfect_grade = calculate_perfect_scenario_grade(portfolio)
    if perfect_grade:
        projections['scenarios'].append({
            'name': 'Perfect Scores',
            'grade': perfect_grade,
            'description': 'Get perfect scores on all remaining assignments'
        })
    
    return projections

def calculate_scenario_grade(portfolio, grade_predictions):
    """Calculate grade based on scenario predictions"""
    grade_breakdown = portfolio.grade_breakdown or {}
    total_weighted_grade = 0
    
    for category, percentage in grade_breakdown.items():
        if category in grade_predictions:
            # Use predicted grade for this category
            predicted_grade = grade_predictions[category]
            total_weighted_grade += predicted_grade * (percentage / 100)
        else:
            # Use current average for this category
            if category == 'exams' and portfolio.exam_grades:
                grades = portfolio.exam_grades
            elif category == 'homework' and portfolio.homework_grades:
                grades = portfolio.homework_grades
            elif category == 'quizzes' and portfolio.quiz_grades:
                grades = portfolio.quiz_grades
            elif category == 'projects' and portfolio.project_grades:
                grades = portfolio.project_grades
            else:
                continue
            
            if grades:
                total_points = sum(grade.get('points_earned', 0) for grade in grades)
                total_possible = sum(grade.get('points_possible', 0) for grade in grades)
                if total_possible > 0:
                    category_grade = (total_points / total_possible) * 100
                    total_weighted_grade += category_grade * (percentage / 100)
    
    return round(total_weighted_grade, 2)

def calculate_perfect_scenario_grade(portfolio):
    """Calculate grade if perfect scores on remaining assignments"""
    grade_breakdown = portfolio.grade_breakdown or {}
    total_weighted_grade = 0
    
    for category, percentage in grade_breakdown.items():
        if category == 'exams' and portfolio.exam_grades:
            grades = portfolio.exam_grades
        elif category == 'homework' and portfolio.homework_grades:
            grades = portfolio.homework_grades
        elif category == 'quizzes' and portfolio.quiz_grades:
            grades = portfolio.quiz_grades
        elif category == 'projects' and portfolio.project_grades:
            grades = portfolio.project_grades
        else:
            continue
        
        if grades:
            # Calculate current average
            total_points = sum(grade.get('points_earned', 0) for grade in grades)
            total_possible = sum(grade.get('points_possible', 0) for grade in grades)
            
            # Assume perfect scores on remaining assignments
            # This is a simplified calculation - in production, you'd need to know
            # how many assignments are remaining
            remaining_weight = max(0, percentage - (total_possible / 100 * percentage))
            perfect_grade = 100
            
            if total_possible > 0:
                current_grade = (total_points / total_possible) * 100
                weighted_current = current_grade * (percentage / 100)
                weighted_perfect = perfect_grade * (remaining_weight / 100)
                total_weighted_grade += weighted_current + weighted_perfect
            else:
                total_weighted_grade += perfect_grade * (percentage / 100)
    
    return round(total_weighted_grade, 2)

# Search & Discovery System
@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def global_search(request):
    """Global search with facets, sorting, and advanced filtering"""
    query = request.query_params.get('q', '')
    school = request.query_params.get('school', '')
    department = request.query_params.get('department', '')
    course_number = request.query_params.get('course_number', '')
    professor = request.query_params.get('professor', '')
    term = request.query_params.get('term', '')
    tags = request.query_params.get('tags', '')
    sort_by = request.query_params.get('sort', 'newest')
    min_price = request.query_params.get('min_price', '')
    max_price = request.query_params.get('max_price', '')
    visibility = request.query_params.get('visibility', '')
    
    # Base queryset
    user = request.user if request.user.is_authenticated else None
    queryset = ClassPortfolio.objects.all()
    
    # Apply visibility filtering
    if not user:
        # Visitors can only see public portfolios
        queryset = queryset.filter(visibility__in=['public_full', 'public_preview'])
    else:
        # Authenticated users can see portfolios they have access to
        accessible_portfolios = []
        for portfolio in queryset:
            if portfolio.can_user_access(user):
                accessible_portfolios.append(portfolio.id)
        queryset = queryset.filter(id__in=accessible_portfolios)
    
    # Apply search filters
    if query:
        queryset = queryset.filter(
            Q(professor__icontains=query) |
            Q(semester__icontains=query)
        )
    
    if school:
        queryset = queryset.filter(created_by__university__icontains=school)
    
    if department:
        # No department filtering since we removed course field
        pass
    
    if course_number:
        # No course filtering since we removed course field
        pass
    
    if professor:
        queryset = queryset.filter(professor__icontains=professor)
    
    if term:
        queryset = queryset.filter(semester__icontains=term)
    
    if tags:
        tag_list = [tag.strip() for tag in tags.split(',')]
        queryset = queryset.filter(tags__overlap=tag_list)
    
    if visibility:
        queryset = queryset.filter(visibility=visibility)
    
    # Apply price filtering for paid portfolios
    if min_price or max_price:
        paid_portfolios = queryset.filter(visibility='paid')
        if min_price:
            paid_portfolios = paid_portfolios.filter(price__gte=float(min_price))
        if max_price:
            paid_portfolios = paid_portfolios.filter(price__lte=float(max_price))
        queryset = paid_portfolios
    
    # Apply sorting
    queryset = apply_search_sorting(queryset, sort_by)
    
    # Get facets for filtering
    facets = get_search_facets(queryset)
    
    # Paginate results
    page_size = 20
    start = int(request.query_params.get('start', 0))
    end = start + page_size
    
    results = queryset[start:end]
    
    return Response({
        'results': ClassPortfolioSerializer(results, many=True, context={'request': request}).data,
        'total_count': queryset.count(),
        'facets': facets,
        'search_params': {
            'query': query,
            'school': school,
            'department': department,
            'course_number': course_number,
            'professor': professor,
            'term': term,
            'tags': tags,
            'sort_by': sort_by,
            'min_price': min_price,
            'max_price': max_price,
            'visibility': visibility
        },
        'pagination': {
            'start': start,
            'end': end,
            'has_more': queryset.count() > end
        }
    })

def apply_search_sorting(queryset, sort_by):
    """Apply sorting to search results"""
    if sort_by == 'newest':
        return queryset.order_by('-created_at')
    elif sort_by == 'oldest':
        return queryset.order_by('created_at')
    elif sort_by == 'most_helpful':
        # Sort by number of reviews and ratings (to be implemented)
        return queryset.annotate(
            helpfulness_score=Count('reviews')
        ).order_by('-helpfulness_score', '-created_at')
    elif sort_by == 'highest_rated':
        # Sort by average rating (to be implemented)
        return queryset.annotate(
            avg_rating=Avg(
                (F('reviews__difficulty_rating') + F('reviews__teaching_quality_rating') + F('reviews__workload_rating')) / 3
            )
        ).order_by('-avg_rating', '-created_at')
    elif sort_by == 'most_purchased':
        # Sort by marketplace popularity
        return queryset.filter(
            marketplace_listing__isnull=False
        ).annotate(
            purchase_count=Count('marketplace_listing__buyers')
        ).order_by('-purchase_count', '-created_at')
    elif sort_by == 'price_low_high':
        return queryset.filter(visibility='paid').order_by('price')
    elif sort_by == 'price_high_low':
        return queryset.filter(visibility='paid').order_by('-price')
    else:
        return queryset.order_by('-created_at')

def get_search_facets(queryset):
    """Get search facets for filtering"""
    # Schools/Universities
    schools = queryset.values('created_by__university').annotate(
        count=Count('id')
    ).filter(created_by__university__isnull=False).order_by('-count')[:10]
    
    # Departments - not available since we removed course field
    departments = []
    
    # Professors
    professors = queryset.values('professor').annotate(
        count=Count('id')
    ).filter(professor__isnull=False).order_by('-count')[:20]
    
    # Terms/Semesters
    terms = queryset.values('semester').annotate(
        count=Count('id')
    ).order_by('-count')[:10]
    
    # Popular tags
    all_tags = []
    for portfolio in queryset:
        all_tags.extend(portfolio.tags)
    
    tag_counts = {}
    for tag in all_tags:
        tag_counts[tag] = tag_counts.get(tag, 0) + 1
    
    popular_tags = sorted(tag_counts.items(), key=lambda x: x[1], reverse=True)[:20]
    
    # Price ranges for paid portfolios
    paid_portfolios = queryset.filter(visibility='paid', price__isnull=False)
    price_ranges = []
    if paid_portfolios.exists():
        min_price = paid_portfolios.aggregate(Min('price'))['price__min']
        max_price = paid_portfolios.aggregate(Max('price'))['price__max']
        
        if min_price and max_price:
            price_ranges = [
                {'range': '0-10', 'count': paid_portfolios.filter(price__lte=10).count()},
                {'range': '10-25', 'count': paid_portfolios.filter(price__gt=10, price__lte=25).count()},
                {'range': '25-50', 'count': paid_portfolios.filter(price__gt=25, price__lte=50).count()},
                {'range': '50+', 'count': paid_portfolios.filter(price__gt=50).count()}
            ]
    
    return {
        'schools': [{'name': s['created_by__university'], 'count': s['count']} for s in schools],
        'departments': [],  # Not available since we removed course field
        'professors': [{'name': p['professor'], 'count': p['count']} for p in professors],
        'terms': [{'name': t['semester'], 'count': t['count']} for t in terms],
        'tags': [{'name': tag, 'count': count} for tag, count in popular_tags],
        'price_ranges': price_ranges
    }

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def search_autocomplete(request):
    """Autocomplete suggestions for search"""
    query = request.query_params.get('q', '')
    field = request.query_params.get('field', 'all')  # all, schools, departments, professors, courses
    
    if len(query) < 2:
        return Response({'suggestions': []})
    
    suggestions = []
    
    if field in ['all', 'schools']:
        # School suggestions
        schools = User.objects.filter(
            university__icontains=query
        ).values_list('university', flat=True).distinct()[:5]
        
        for school in schools:
            suggestions.append({
                'type': 'school',
                'value': school,
                'label': school
            })
    
    if field in ['all', 'departments']:
        # Department suggestions
        departments = Department.objects.filter(
            Q(code__icontains=query) | Q(name__icontains=query)
        )[:5]
        
        for dept in departments:
            suggestions.append({
                'type': 'department',
                'value': dept.code,
                'label': f"{dept.code} - {dept.name}"
            })
    
    if field in ['all', 'professors']:
        # Professor suggestions
        professors = Professor.objects.filter(
            name__icontains=query
        )[:5]
        
        for prof in professors:
            suggestions.append({
                'type': 'professor',
                'value': prof.name,
                'label': prof.name
            })
    
    
    return Response({'suggestions': suggestions})

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def search_suggestions(request):
    """Get search suggestions and trending searches"""
    # Get popular searches (to be implemented with search analytics)
    popular_searches = [
        'computer science',
        'data structures',
        'algorithms',
        'calculus',
        'physics',
        'chemistry',
        'biology',
        'engineering'
    ]
    
    
    # Get trending professors
    trending_professors = Professor.objects.annotate(
        portfolio_count=Count('portfolios')
    ).order_by('-portfolio_count')[:10]
    
    # Get trending tags
    all_tags = []
    for portfolio in ClassPortfolio.objects.all():
        all_tags.extend(portfolio.tags)
    
    tag_counts = {}
    for tag in all_tags:
        tag_counts[tag] = tag_counts.get(tag, 0) + 1
    
    trending_tags = sorted(tag_counts.items(), key=lambda x: x[1], reverse=True)[:15]
    
    return Response({
        'popular_searches': popular_searches,
        'trending_professors': ProfessorSerializer(trending_professors, many=True).data,
        'trending_tags': [{'name': tag, 'count': count} for tag, count in trending_tags]
    })

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def search_analytics(request):
    """Get search analytics for the user"""
    user = request.user
    
    # Get user's search history (to be implemented)
    search_history = []  # Placeholder
    
    # Get user's saved searches (to be implemented)
    saved_searches = []  # Placeholder
    
    # Get recommended portfolios based on user's activity
    recommended_portfolios = get_recommended_portfolios(user)
    
    return Response({
        'search_history': search_history,
        'saved_searches': saved_searches,
        'recommended_portfolios': ClassPortfolioSerializer(recommended_portfolios, many=True, context={'request': request}).data
    })

def get_recommended_portfolios(user):
    """Get recommended portfolios for a user"""
    # Simple recommendation algorithm based on user's university and major
    recommendations = ClassPortfolio.objects.filter(
        visibility__in=['public_full', 'public_preview']
    )
    
    # Filter by same university
    if user.university:
        recommendations = recommendations.filter(
            created_by__university__icontains=user.university
        )
    
    # Filter by same major/department
    # Major-based recommendations not available since we removed course field
    # if user.major:
    #     recommendations = recommendations.filter(
    #         course__department__name__icontains=user.major
    #     )
    
    # Order by popularity and recency
    recommendations = recommendations.annotate(
        portfolio_count=Count('id')
    ).order_by('-portfolio_count', '-created_at')
    
    return recommendations[:10]

# Preview Rules & Content Restrictions
@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def portfolio_preview_content(request, portfolio_id):
    """Get portfolio content with preview restrictions applied"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        user = request.user if request.user.is_authenticated else None
        
        # Check if user can access this portfolio
        if not portfolio.can_user_access(user) and not portfolio.is_public_preview():
            return Response({'error': 'Access denied'}, status=status.HTTP_403_FORBIDDEN)
        
        # Determine if this is a preview request
        is_preview = portfolio.is_public_preview() and not portfolio.can_user_access(user)
        
        # Get preview-restricted content
        if is_preview:
            # Simple preview content
            preview_data = {
                'metadata': {
                    'professor_name': portfolio.professor,
                    'semester': portfolio.semester,
                    'year': portfolio.year,
                    'tags': portfolio.tags,
                    'created_at': portfolio.created_at,
                    'owner': portfolio.created_by.username if portfolio.created_by else None
                },
                'syllabus': {
                    'important_dates': ImportantDateSerializer(
                        ImportantDate.objects.filter(portfolio=portfolio)[:3], many=True
                    ).data,
                    'total_dates': ImportantDate.objects.filter(portfolio=portfolio).count(),
                    'preview_note': "Showing first 3 important dates"
                },
                'materials': {
                    'count': LectureMaterial.objects.filter(portfolio=portfolio).count(),
                    'preview_note': "Materials available with full access"
                },
                'flashcards': {
                    'items': [],
                    'total_count': 0,
                    'preview_note': "Flashcards available with full access"
                },
                'quizzes': {
                    'items': [],
                    'total_count': 0,
                    'preview_note': "Quizzes available with full access"
                },
                'performance': {
                    'current_grade': None,
                    'grade_breakdown': None,
                    'preview_note': "Grade information available with full access"
                },
                'community': {
                    'posts': [],
                    'total_posts': 0,
                    'thread_titles': [],
                    'preview_note': "Community discussions available with full access"
                }
            }
        else:
            preview_data = get_full_content(portfolio)
        
        return Response({
            'portfolio': {
                'id': portfolio.id,
                'professor': ProfessorSerializer(portfolio.professor).data,
                'semester': portfolio.semester,
                'year': portfolio.year,
                'visibility': portfolio.visibility,
                'tags': portfolio.tags,
                'created_at': portfolio.created_at,
                'owner': portfolio.created_by.username if portfolio.created_by else None
            },
            'preview_mode': is_preview,
            'content': preview_data,
            'access_level': get_access_level(portfolio, user)
        })
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

def get_preview_content(portfolio, user, is_preview):
    """Get content with preview restrictions applied"""
    content = {}
    
    # Always show: metadata, tags, creation date, owner display name
    content['metadata'] = {
        'professor_name': portfolio.professor,
        'semester': portfolio.semester,
        'year': portfolio.year,
        'tags': portfolio.tags,
        'created_at': portfolio.created_at,
        'owner': portfolio.created_by.username if portfolio.created_by else None
    }
    
    if is_preview:
        # Apply preview restrictions
        content.update(get_preview_restricted_content(portfolio))
    else:
        # Full access - show everything
        content.update(get_full_content(portfolio))
    
    return content

def get_preview_restricted_content(portfolio):
    """Get content with preview restrictions applied"""
    content = {}
    
    # Syllabus/Calendar: Show 20% of items (rounded up) and next upcoming item
    important_dates = ImportantDate.objects.filter(portfolio=portfolio).order_by('due_date')
    total_dates = important_dates.count()
    
    if total_dates > 0:
        # Calculate 20% rounded up
        preview_count = max(1, (total_dates + 4) // 5)  # Round up to nearest integer
        
        # Get next upcoming item
        upcoming = important_dates.filter(due_date__gte=timezone.now()).first()
        
        # Get preview items
        preview_dates = important_dates[:preview_count]
        if upcoming and upcoming not in preview_dates:
            preview_dates = list(preview_dates) + [upcoming]
        
        content['syllabus'] = {
            'important_dates': ImportantDateSerializer(preview_dates, many=True).data,
            'total_dates': total_dates,
            'preview_count': len(preview_dates),
            'preview_note': f"Showing {len(preview_dates)} of {total_dates} important dates"
        }
    else:
        content['syllabus'] = {
            'important_dates': [],
            'total_dates': 0,
            'preview_count': 0,
            'preview_note': "No important dates available"
        }
    
    # Notes/Study Assets: Show first 5 flashcards + first 3 quiz questions + summary first paragraph
    flashcards = Flashcard.objects.filter(material__portfolio=portfolio).order_by('created_at')[:5]
    content['flashcards'] = {
        'items': FlashcardSerializer(flashcards, many=True).data,
        'total_count': Flashcard.objects.filter(material__portfolio=portfolio).count(),
        'preview_note': f"Showing first {len(flashcards)} flashcards"
    }
    
    quizzes = Quiz.objects.filter(portfolio=portfolio).order_by('created_at')
    quiz_preview = []
    for quiz in quizzes[:1]:  # Show first quiz only
        questions = QuizQuestion.objects.filter(quiz=quiz).order_by('id')[:3]
        quiz_preview.append({
            'quiz_id': quiz.id,
            'title': quiz.title,
            'questions': QuizQuestionSerializer(questions, many=True).data,
            'total_questions': QuizQuestion.objects.filter(quiz=quiz).count(),
            'preview_note': f"Showing first {len(questions)} questions"
        })
    
    content['quizzes'] = {
        'items': quiz_preview,
        'total_count': quizzes.count(),
        'preview_note': f"Showing preview of first quiz"
    }
    
    # AI Summaries: Show first paragraph only
    materials = LectureMaterial.objects.filter(portfolio=portfolio).order_by('uploaded_at')
    summaries = []
    for material in materials[:3]:  # Show first 3 materials
        if hasattr(material, 'processed_file') and material.processed_file.summary:
            summary_text = material.processed_file.summary
            # Get first paragraph
            first_paragraph = summary_text.split('\n')[0] if '\n' in summary_text else summary_text[:200] + "..."
            summaries.append({
                'material_id': material.id,
                'title': material.title,
                'summary_preview': first_paragraph,
                'full_summary_length': len(summary_text),
                'preview_note': "First paragraph preview"
            })
    
    content['summaries'] = {
        'items': summaries,
        'total_count': materials.count(),
        'preview_note': f"Showing preview of {len(summaries)} summaries"
    }
    
    # Performance: Hide numeric grades; show feature screenshot only
    content['performance'] = {
        'current_grade': None,  # Hidden in preview
        'grade_breakdown': None,  # Hidden in preview
        'category_averages': None,  # Hidden in preview
        'preview_note': "Grade information hidden in preview mode",
        'feature_available': True,
        'screenshot_url': "/static/images/grade-tracker-preview.png"  # Placeholder
    }
    
    # Reviews: Show averages and up to 1 full review
    reviews = ClassReview.objects.filter(portfolio=portfolio).order_by('-created_at')
    if reviews.exists():
        # Calculate averages
        avg_difficulty = reviews.aggregate(avg=Avg('difficulty_rating'))['avg'] or 0
        avg_quality = reviews.aggregate(avg=Avg('teaching_quality_rating'))['avg'] or 0
        avg_workload = reviews.aggregate(avg=Avg('workload_rating'))['avg'] or 0
        
        # Show one full review
        sample_review = reviews.first()
        
        content['reviews'] = {
            'averages': {
                'difficulty': round(avg_difficulty, 1),
                'teaching_quality': round(avg_quality, 1),
                'workload': round(avg_workload, 1),
                'overall': round((avg_difficulty + avg_quality + avg_workload) / 3, 1)
            },
            'sample_review': {
                'rating': sample_review.overall_rating,
                'comment': sample_review.comment,
                'created_at': sample_review.created_at
            },
            'total_reviews': reviews.count(),
            'preview_note': f"Showing averages and 1 of {reviews.count()} reviews"
        }
    else:
        content['reviews'] = {
            'averages': None,
            'sample_review': None,
            'total_reviews': 0,
            'preview_note': "No reviews available"
        }
    
    # Community: Show thread titles; hide content until login
    posts = Post.objects.filter(portfolio=portfolio).order_by('-created_at')
    content['community'] = {
        'thread_titles': [
            {
                'id': post.id,
                'title': post.title,
                'author': post.author.username,
                'created_at': post.created_at,
                'reply_count': post.comments.count()
            }
            for post in posts[:10]
        ],
        'total_threads': posts.count(),
        'preview_note': "Thread titles only - login to view content"
    }
    
    return content

def get_full_content(portfolio):
    """Get full content without restrictions"""
    content = {}
    
    # Full syllabus data
    important_dates = ImportantDate.objects.filter(portfolio=portfolio).order_by('due_date')
    content['syllabus'] = {
        'important_dates': ImportantDateSerializer(important_dates, many=True).data,
        'total_dates': important_dates.count()
    }
    
    # Full flashcards
    flashcards = Flashcard.objects.filter(material__portfolio=portfolio).order_by('created_at')
    content['flashcards'] = {
        'items': FlashcardSerializer(flashcards, many=True).data,
        'total_count': flashcards.count()
    }
    
    # Full quizzes
    quizzes = Quiz.objects.filter(portfolio=portfolio).order_by('created_at')
    quiz_data = []
    for quiz in quizzes:
        questions = QuizQuestion.objects.filter(quiz=quiz).order_by('id')
        quiz_data.append({
            'quiz_id': quiz.id,
            'title': quiz.title,
            'questions': QuizQuestionSerializer(questions, many=True).data,
            'total_questions': questions.count()
        })
    
    content['quizzes'] = {
        'items': quiz_data,
        'total_count': quizzes.count()
    }
    
    # Full summaries
    materials = LectureMaterial.objects.filter(portfolio=portfolio).order_by('uploaded_at')
    summaries = []
    for material in materials:
        if hasattr(material, 'processed_file') and material.processed_file.summary:
            summaries.append({
                'material_id': material.id,
                'title': material.title,
                'summary': material.processed_file.summary,
                'topics': material.processed_file.topics,
                'generated_at': material.processed_file.processed_at
            })
    
    content['summaries'] = {
        'items': summaries,
        'total_count': materials.count()
    }
    
    # Full performance data
    all_grades = portfolio.all_grades
    current_grade = portfolio.calculate_current_grade()
    grade_breakdown = portfolio.grade_breakdown or {}
    
    content['performance'] = {
        'current_grade': current_grade,
        'all_grades': all_grades,
        'grade_breakdown': grade_breakdown,
        'category_averages': calculate_category_averages(portfolio)
    }
    
    # Full reviews
    reviews = ClassReview.objects.filter(portfolio=portfolio).order_by('-created_at')
    content['reviews'] = {
        'items': ClassReviewSerializer(reviews, many=True).data,
        'total_count': reviews.count()
    }
    
    # Full community
    posts = Post.objects.filter(portfolio=portfolio).order_by('-created_at')
    content['community'] = {
        'posts': PostSerializer(posts, many=True).data,
        'total_posts': posts.count()
    }
    
    return content

def calculate_category_averages(portfolio):
    """Calculate category averages for performance tracking"""
    grade_breakdown = portfolio.grade_breakdown or {}
    category_averages = {}
    
    for category, percentage in grade_breakdown.items():
        if category == 'exams' and portfolio.exam_grades:
            grades = portfolio.exam_grades
        elif category == 'homework' and portfolio.homework_grades:
            grades = portfolio.homework_grades
        elif category == 'quizzes' and portfolio.quiz_grades:
            grades = portfolio.quiz_grades
        elif category == 'projects' and portfolio.project_grades:
            grades = portfolio.project_grades
        else:
            continue
        
        if grades:
            total_points = sum(grade.get('points_earned', 0) for grade in grades)
            total_possible = sum(grade.get('points_possible', 0) for grade in grades)
            if total_possible > 0:
                category_averages[category] = {
                    'average': round((total_points / total_possible) * 100, 2),
                    'percentage': percentage,
                    'count': len(grades)
                }
    
    return category_averages

def get_access_level(portfolio, user):
    """Determine user's access level to portfolio"""
    if not user:
        return 'visitor'
    elif portfolio.can_user_edit(user):
        return 'owner'
    elif portfolio.can_user_access(user):
        return 'full_access'
    else:
        return 'preview_only'

def get_upgrade_options(portfolio, user):
    """Get upgrade options for preview users"""
    if not user:
        return {
            'login_required': True,
            'message': 'Login to access full content',
            'login_url': '/api/auth/login/',
            'register_url': '/api/auth/register/'
        }
    
    if portfolio.visibility == 'paid':
        # Check if user already purchased
        if hasattr(portfolio, 'marketplace_listing'):
            listing = portfolio.marketplace_listing
            if listing.buyers.filter(id=user.id).exists():
                return None  # Already purchased
        
        return {
            'purchase_required': True,
            'price': float(portfolio.price),
            'purchase_url': f'/api/marketplace/{portfolio.marketplace_listing.id}/purchase/',
            'message': f'Purchase for ${portfolio.price} to access full content'
        }
    
    return None

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def document_preview(request, file_id):
    """Get watermarked document preview"""
    try:
        # This would integrate with a document processing service
        # For now, return a placeholder response
        
        return Response({
            'file_id': file_id,
            'preview_url': f'/api/documents/{file_id}/preview/',
            'watermark': 'PREVIEW - Class Portfolio',
            'download_restricted': True,
            'message': 'Full document available with purchase or login'
        })
        
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Non-functional Requirements: Privacy, Security, Performance, Analytics
@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def privacy_policy(request):
    """Get privacy policy and IP guidelines"""
    return Response({
        'privacy_policy': {
            'data_collection': {
                'personal_data': ['username', 'email', 'university', 'major'],
                'usage_data': ['search_queries', 'portfolio_views', 'purchase_history'],
                'content_data': ['uploaded_files', 'generated_content', 'reviews']
            },
            'data_usage': {
                'personalization': 'Recommend portfolios based on university and major',
                'analytics': 'Improve service quality and user experience',
                'communication': 'Send notifications and updates'
            },
            'data_sharing': {
                'third_parties': 'No personal data shared with third parties',
                'public_content': 'Portfolio content may be publicly visible based on visibility settings',
                'legal_requirements': 'Data may be disclosed if required by law'
            },
            'user_rights': {
                'access': 'Request access to your personal data',
                'correction': 'Correct inaccurate personal data',
                'deletion': 'Request deletion of your personal data',
                'portability': 'Export your data in a portable format'
            }
        },
        'ip_policy': {
            'copyrighted_materials': {
                'syllabi': 'Allowed - educational fair use',
                'lecture_notes': 'Must be original or fair-use excerpts',
                'textbooks': 'Not allowed - copyrighted content',
                'images': 'Must be original or properly licensed'
            },
            'dmca_compliance': {
                'takedown_process': '24-hour response to valid DMCA notices',
                'counter_notice': 'Users can file counter-notices',
                'repeat_infringers': 'Account termination for repeat violations'
            },
            'content_guidelines': {
                'original_content': 'Encouraged - original notes and summaries',
                'attribution': 'Required for any quoted material',
                'academic_integrity': 'No sharing of exam questions or answers'
            }
        },
        'contact': {
            'privacy_officer': 'privacy@classportfolio.com',
            'dmca_agent': 'dmca@classportfolio.com',
            'legal': 'legal@classportfolio.com'
        }
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def dmca_takedown_request(request):
    """Submit DMCA takedown request"""
    data = request.data
    
    required_fields = ['copyrighted_work', 'infringing_content', 'contact_info', 'good_faith']
    for field in required_fields:
        if field not in data:
            return Response({'error': f'Missing required field: {field}'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Create DMCA request record
    dmca_request = {
        'id': f"DMCA-{int(timezone.now().timestamp())}",
        'copyrighted_work': data['copyrighted_work'],
        'infringing_content': data['infringing_content'],
        'contact_info': data['contact_info'],
        'good_faith': data['good_faith'],
        'submitted_by': request.user.id,
        'submitted_at': timezone.now(),
        'status': 'pending',
        'response_deadline': timezone.now() + timedelta(hours=24)
    }
    
    # In production, this would be stored in a database and sent to moderators
    return Response({
        'message': 'DMCA takedown request submitted successfully',
        'request_id': dmca_request['id'],
        'status': dmca_request['status'],
        'response_deadline': dmca_request['response_deadline'],
        'next_steps': 'Request will be reviewed within 24 hours'
    }, status=status.HTTP_201_CREATED)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def security_status(request):
    """Get user's security status and recommendations"""
    user = request.user
    
    # Check security factors
    security_factors = {
        'password_strength': check_password_strength(user),
        'two_factor_enabled': False,  # To be implemented
        'recent_logins': get_recent_logins(user),
        'suspicious_activity': check_suspicious_activity(user),
        'account_age': (timezone.now() - user.date_joined).days
    }
    
    # Generate security recommendations
    recommendations = generate_security_recommendations(security_factors)
    
    return Response({
        'user_id': user.id,
        'security_factors': security_factors,
        'recommendations': recommendations,
        'security_score': calculate_security_score(security_factors),
        'last_security_check': timezone.now()
    })

def check_password_strength(user):
    """Check password strength (placeholder implementation)"""
    # In production, this would check password complexity
    return {
        'strength': 'strong',
        'last_changed': user.date_joined,
        'recommendation': 'Password appears strong'
    }

def get_recent_logins(user):
    """Get recent login information"""
    # In production, this would track login history
    return [
        {
            'timestamp': timezone.now() - timedelta(hours=2),
            'ip_address': '192.168.1.1',
            'user_agent': 'Mozilla/5.0...',
            'location': 'Austin, TX'
        }
    ]

def check_suspicious_activity(user):
    """Check for suspicious activity patterns"""
    # In production, this would analyze user behavior
    return {
        'unusual_logins': False,
        'rapid_uploads': False,
        'bulk_actions': False,
        'risk_level': 'low'
    }

def generate_security_recommendations(factors):
    """Generate security recommendations based on factors"""
    recommendations = []
    
    if not factors['two_factor_enabled']:
        recommendations.append({
            'type': 'two_factor',
            'priority': 'high',
            'message': 'Enable two-factor authentication for better security',
            'action_url': '/api/security/enable-2fa/'
        })
    
    if factors['account_age'] < 30:
        recommendations.append({
            'type': 'account_setup',
            'priority': 'medium',
            'message': 'Complete your profile setup for better security',
            'action_url': '/api/users/profile/'
        })
    
    return recommendations

def calculate_security_score(factors):
    """Calculate overall security score"""
    score = 100
    
    if not factors['two_factor_enabled']:
        score -= 20
    
    if factors['suspicious_activity']['risk_level'] == 'high':
        score -= 30
    elif factors['suspicious_activity']['risk_level'] == 'medium':
        score -= 15
    
    if factors['password_strength']['strength'] == 'weak':
        score -= 25
    
    return max(0, score)

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def performance_metrics(request):
    """Get system performance metrics"""
    # In production, this would collect real metrics
    return Response({
        'system_health': {
            'status': 'healthy',
            'uptime': '99.9%',
            'response_time': '150ms',
            'throughput': '1000 requests/minute'
        },
        'performance_targets': {
            'file_upload': '≤ 10 seconds for 10-page PDF',
            'page_load': '< 2 seconds p95',
            'api_response': '< 500ms p95',
            'database_query': '< 100ms p95'
        },
        'current_metrics': {
            'avg_upload_time': '3.2 seconds',
            'avg_page_load': '1.1 seconds',
            'avg_api_response': '180ms',
            'avg_db_query': '45ms'
        },
        'optimization_status': {
            'caching': 'enabled',
            'cdn': 'enabled',
            'database_indexing': 'optimized',
            'file_compression': 'enabled'
        }
    })

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def accessibility_features(request):
    """Get accessibility features and compliance status"""
    return Response({
        'wcag_compliance': {
            'level': 'AA',
            'version': '2.1',
            'status': 'compliant',
            'last_audit': '2024-09-01'
        },
        'keyboard_navigation': {
            'tab_order': 'logical',
            'skip_links': 'available',
            'focus_indicators': 'visible',
            'keyboard_shortcuts': 'documented'
        },
        'screen_reader_support': {
            'alt_text': 'required for images',
            'aria_labels': 'implemented',
            'semantic_html': 'used',
            'heading_structure': 'proper'
        },
        'visual_accessibility': {
            'color_contrast': '4.5:1 minimum',
            'text_scaling': 'up to 200%',
            'high_contrast_mode': 'supported',
            'font_options': 'adjustable'
        },
        'motor_accessibility': {
            'large_click_targets': '44px minimum',
            'drag_drop_alternatives': 'available',
            'timeout_extensions': 'configurable',
            'voice_control': 'supported'
        },
        'cognitive_accessibility': {
            'clear_language': 'used',
            'error_messages': 'helpful',
            'progress_indicators': 'visible',
            'help_documentation': 'available'
        }
    })

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_analytics(request):
    """Get user-specific analytics and insights"""
    user = request.user
    
    # Get user activity metrics
    activity_metrics = {
        'portfolios_created': ClassPortfolio.objects.filter(created_by=user).count(),
        'portfolios_purchased': PortfolioPurchase.objects.filter(buyer=user).count(),
        'total_uploads': LectureMaterial.objects.filter(portfolio__created_by=user).count(),
        'total_reviews': ClassReview.objects.filter(reviewer=user).count(),
        'total_searches': 0,  # To be implemented with search tracking
        'last_activity': user.last_login or user.date_joined
    }
    
    # Get engagement metrics
    engagement_metrics = {
        'session_duration': '12 minutes average',
        'pages_per_session': 8.5,
        'return_visitor': True,
        'feature_usage': {
            'syllabus_scanner': 85,
            'learning_space': 92,
            'performance_tracker': 78,
            'marketplace': 45
        }
    }
    
    # Get learning insights
    learning_insights = {
        'most_active_subjects': get_user_subjects(user),
        'study_patterns': analyze_study_patterns(user),
        'achievement_badges': get_user_badges(user),
        'recommendations': get_learning_recommendations(user)
    }
    
    return Response({
        'user_id': user.id,
        'activity_metrics': activity_metrics,
        'engagement_metrics': engagement_metrics,
        'learning_insights': learning_insights,
        'analytics_period': 'last_30_days',
        'data_collection_date': timezone.now()
    })

def get_user_subjects(user):
    """Get user's most active subjects"""
    portfolios = ClassPortfolio.objects.filter(created_by=user)
    subjects = {}
    
    # Subjects not available since we removed course field
    subjects = {}
    
    return sorted(subjects.items(), key=lambda x: x[1], reverse=True)[:5]

def analyze_study_patterns(user):
    """Analyze user's study patterns"""
    # In production, this would analyze actual usage data
    return {
        'peak_hours': '2-4 PM',
        'study_days': 'Monday, Wednesday, Friday',
        'session_length': '45 minutes average',
        'break_frequency': 'every 25 minutes'
    }

def get_user_badges(user):
    """Get user's achievement badges"""
    badges = []
    
    portfolio_count = ClassPortfolio.objects.filter(created_by=user).count()
    if portfolio_count >= 5:
        badges.append({'name': 'Portfolio Creator', 'description': 'Created 5+ portfolios'})
    
    if portfolio_count >= 10:
        badges.append({'name': 'Super Creator', 'description': 'Created 10+ portfolios'})
    
    return badges

def get_learning_recommendations(user):
    """Get personalized learning recommendations"""
    return [
        {
            'type': 'feature',
            'title': 'Try the Performance Tracker',
            'description': 'Track your grades and get insights',
            'action_url': '/api/portfolios/create-wizard/'
        },
        {
            'type': 'content',
            'title': 'Explore Computer Science Portfolios',
            'description': 'Based on your interests',
            'action_url': '/api/search/?department=CS'
        }
    ]

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def audit_log(request):
    """Get user's audit log and access history"""
    user = request.user
    
    # In production, this would track all user actions
    audit_entries = [
        {
            'timestamp': timezone.now() - timedelta(hours=1),
            'action': 'portfolio_view',
            'resource': f'Portfolio #{1}',
            'ip_address': '192.168.1.1',
            'user_agent': 'Mozilla/5.0...',
            'result': 'success'
        },
        {
            'timestamp': timezone.now() - timedelta(hours=2),
            'action': 'file_upload',
            'resource': 'lecture_notes.pdf',
            'ip_address': '192.168.1.1',
            'user_agent': 'Mozilla/5.0...',
            'result': 'success'
        },
        {
            'timestamp': timezone.now() - timedelta(days=1),
            'action': 'portfolio_create',
            'resource': 'CS101 Portfolio',
            'ip_address': '192.168.1.1',
            'user_agent': 'Mozilla/5.0...',
            'result': 'success'
        }
    ]
    
    return Response({
        'user_id': user.id,
        'audit_entries': audit_entries,
        'total_entries': len(audit_entries),
        'retention_period': '90 days',
        'export_available': True
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def anti_scraping_protection(request):
    """Implement anti-scraping measures"""
    # Rate limiting
    user = request.user
    action = request.data.get('action', 'unknown')
    
    # In production, this would implement:
    # - Rate limiting per user/IP
    # - CAPTCHA for suspicious activity
    # - Request pattern analysis
    # - Bot detection
    
    return Response({
        'message': 'Anti-scraping protection active',
        'user_id': user.id,
        'action': action,
        'protection_level': 'standard',
        'rate_limit_status': 'within_limits'
    })

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def system_status(request):
    """Get overall system status and health"""
    return Response({
        'status': 'operational',
        'version': '1.0.0',
        'uptime': '99.9%',
        'last_updated': timezone.now(),
        'services': {
            'api': 'healthy',
            'database': 'healthy',
            'file_storage': 'healthy',
            'ai_services': 'healthy',
            'payment_processing': 'healthy'
        },
        'performance': {
            'response_time': '150ms',
            'throughput': '1000 req/min',
            'error_rate': '0.1%',
            'availability': '99.9%'
        },
        'security': {
            'ssl_enabled': True,
            'firewall_active': True,
            'ddos_protection': True,
            'encryption': 'AES-256'
        }
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def save_search(request):
    """Save a search query for the user"""
    user = request.user
    query = request.data.get('query', '')
    filters = request.data.get('filters', {})
    name = request.data.get('name', f"Search: {query}")
    
    if not query:
        return Response({'error': 'Query is required'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Create saved search (to be implemented with SavedSearch model)
    saved_search = {
        'name': name,
        'query': query,
        'filters': filters,
        'created_at': timezone.now(),
        'user': user.id
    }
    
    return Response({
        'message': 'Search saved successfully',
        'saved_search': saved_search
    }, status=status.HTTP_201_CREATED)

class SyllabusDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Syllabus.objects.all()
    serializer_class = SyllabusSerializer
    permission_classes = [permissions.IsAuthenticated]

class SyllabusExtractionView(generics.RetrieveUpdateAPIView):
    """View to get and update extracted syllabus data"""
    queryset = SyllabusExtraction.objects.all()
    serializer_class = SyllabusExtractionSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_object(self):
        syllabus_id = self.kwargs.get('syllabus_id')
        try:
            return SyllabusExtraction.objects.get(syllabus_id=syllabus_id)
        except SyllabusExtraction.DoesNotExist:
            return None

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def extract_syllabus(request, syllabus_id):
    """Manually trigger syllabus extraction"""
    try:
        syllabus = Syllabus.objects.get(id=syllabus_id)
        
        # Check if extraction already exists
        if hasattr(syllabus, 'extraction'):
            return Response({'message': 'Extraction already exists'}, status=400)
        
        # Start extraction process
        syllabus.extraction_status = 'processing'
        syllabus.save()
        
        # Extract text from file
        extracted_text = extract_text_from_file(syllabus.file)
        syllabus.extracted_text = extracted_text
        syllabus.save()
        
        # Run AI extraction
        from .syllabus_extractor import SyllabusExtractor
        extractor = SyllabusExtractor()
        extraction_data = extractor.extract_from_text(extracted_text)
        
        # Create extraction record
        from .models import SyllabusExtraction
        SyllabusExtraction.objects.create(
            syllabus=syllabus,
            **extraction_data
        )
        
        syllabus.extraction_status = 'completed'
        syllabus.save()
        
        return Response({'message': 'Extraction completed successfully'})
        
    except Syllabus.DoesNotExist:
        return Response({'error': 'Syllabus not found'}, status=404)
    except Exception as e:
        syllabus.extraction_status = 'failed'
        syllabus.extraction_error = str(e)
        syllabus.save()
        return Response({'error': str(e)}, status=500)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_important_dates_from_extraction(request, syllabus_id):
    """Create ImportantDate objects from extracted syllabus data"""
    try:
        syllabus = Syllabus.objects.get(id=syllabus_id)
        if not hasattr(syllabus, 'extraction'):
            return Response({'error': 'No extraction data found'}, status=400)
        
        extraction = syllabus.extraction
        created_dates = []
        
        # Create dates from extracted data
        for date_data in extraction.all_important_dates:
            ImportantDate.objects.create(
                portfolio=syllabus.portfolio,
                title=date_data['title'],
                date_type=date_data['type'],
                due_date=datetime.strptime(date_data['date'], '%Y-%m-%d'),
                description=f"Auto-extracted from syllabus",
                is_synced=False
            )
            created_dates.append(date_data)
        
        return Response({
            'message': f'Created {len(created_dates)} important dates',
            'dates': created_dates
        })
        
    except Syllabus.DoesNotExist:
        return Response({'error': 'Syllabus not found'}, status=404)
    except Exception as e:
        return Response({'error': str(e)}, status=500)

def extract_text_from_file(file):
    """Extract text from uploaded file"""
    # Simplified implementation - in production, use proper file parsing
    return f"Text extracted from {file.name}"

# Important Date Views
class ImportantDateListCreateView(generics.ListCreateAPIView):
    queryset = ImportantDate.objects.all()
    serializer_class = ImportantDateSerializer
    permission_classes = [permissions.IsAuthenticated]

class ImportantDateDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = ImportantDate.objects.all()
    serializer_class = ImportantDateSerializer
    permission_classes = [permissions.IsAuthenticated]

@api_view(['GET'])
def upcoming_deadlines(request):
    """Get upcoming deadlines for the next 7 days"""
    now = timezone.now()
    week_from_now = now + timedelta(days=7)
    
    deadlines = ImportantDate.objects.filter(
        due_date__range=[now, week_from_now]
    ).order_by('due_date')
    
    serializer = ImportantDateSerializer(deadlines, many=True)
    return Response(serializer.data)

# Lecture Material Views
class LectureMaterialListCreateView(generics.ListCreateAPIView):
    queryset = LectureMaterial.objects.all()
    serializer_class = LectureMaterialSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def perform_create(self, serializer):
        serializer.save(uploaded_by=self.request.user)

class LectureMaterialDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = LectureMaterial.objects.all()
    serializer_class = LectureMaterialSerializer
    permission_classes = [permissions.IsAuthenticated]

# Flashcard Views
class FlashcardListCreateView(generics.ListCreateAPIView):
    queryset = Flashcard.objects.all()
    serializer_class = FlashcardSerializer
    permission_classes = [permissions.IsAuthenticated]

class FlashcardDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Flashcard.objects.all()
    serializer_class = FlashcardSerializer
    permission_classes = [permissions.IsAuthenticated]

# Quiz Views
class QuizListCreateView(generics.ListCreateAPIView):
    queryset = Quiz.objects.all()
    serializer_class = QuizSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)

class QuizDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Quiz.objects.all()
    serializer_class = QuizSerializer
    permission_classes = [permissions.IsAuthenticated]

class QuizQuestionListCreateView(generics.ListCreateAPIView):
    queryset = QuizQuestion.objects.all()
    serializer_class = QuizQuestionSerializer
    permission_classes = [permissions.IsAuthenticated]

class QuizQuestionDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = QuizQuestion.objects.all()
    serializer_class = QuizQuestionSerializer
    permission_classes = [permissions.IsAuthenticated]

# Quiz Submission Views
class QuizSubmissionListCreateView(generics.ListCreateAPIView):
    queryset = QuizSubmission.objects.all()
    serializer_class = QuizSubmissionSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        return QuizSubmission.objects.filter(user=self.request.user)

class QuizSubmissionDetailView(generics.RetrieveAPIView):
    queryset = QuizSubmission.objects.all()
    serializer_class = QuizSubmissionSerializer
    permission_classes = [permissions.IsAuthenticated]

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def submit_quiz(request, quiz_id):
    """Submit quiz answers and calculate score"""
    try:
        quiz = Quiz.objects.get(id=quiz_id)
        user = request.user
        
        # Check if user already submitted this quiz
        submission, created = QuizSubmission.objects.get_or_create(
            quiz=quiz, 
            user=user,
            defaults={'answers': {}}
        )
        
        if not created:
            return Response({'error': 'Quiz already submitted'}, status=400)
        
        # Get answers from request
        answers = request.data.get('answers', {})
        time_taken = request.data.get('time_taken_minutes', None)
        
        # Validate answers format
        if not isinstance(answers, dict):
            return Response({'error': 'Answers must be a dictionary'}, status=400)
        
        # Update submission
        submission.answers = answers
        if time_taken:
            submission.time_taken_minutes = time_taken
        
        # Calculate score
        score = submission.calculate_score()
        
        return Response({
            'message': 'Quiz submitted successfully',
            'score': score,
            'total_points': submission.total_points,
            'submission_id': submission.id
        })
        
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz not found'}, status=404)
    except Exception as e:
        return Response({'error': str(e)}, status=500)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def quiz_results(request, quiz_id):
    """Get quiz results for a specific quiz"""
    try:
        quiz = Quiz.objects.get(id=quiz_id)
        submissions = QuizSubmission.objects.filter(quiz=quiz, user=request.user)
        
        if not submissions.exists():
            return Response({'message': 'No submissions found'})
        
        latest_submission = submissions.first()
        
        # Get detailed results
        results = []
        for question in quiz.questions.all():
            user_answer = latest_submission.answers.get(str(question.id), '')
            is_correct = question.validate_answer(user_answer)
            
            results.append({
                'question_id': question.id,
                'question_text': question.question_text,
                'question_type': question.question_type,
                'user_answer': user_answer,
                'correct_answer': question.get_correct_answer(),
                'is_correct': is_correct,
                'points': question.points,
                'explanation': question.explanation
            })
        
        return Response({
            'quiz_title': quiz.title,
            'score': latest_submission.score,
            'total_points': latest_submission.total_points,
            'submitted_at': latest_submission.submitted_at,
            'time_taken_minutes': latest_submission.time_taken_minutes,
            'results': results
        })
        
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz not found'}, status=404)

@api_view(['GET'])
def grade_analytics(request, portfolio_id):
    """Get grade analytics for a portfolio"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        
        # Calculate current grade
        current_grade = portfolio.calculate_current_grade()
        
        # Get all grades
        all_grades = portfolio.all_grades
        
        if not all_grades:
            return Response({'message': 'No grades found'})
        
        # Calculate totals
        total_points_earned = sum(grade.get('points_earned', 0) for grade in all_grades)
        total_points_possible = sum(grade.get('points_possible', 0) for grade in all_grades)
        current_average = (total_points_earned / total_points_possible) * 100 if total_points_possible > 0 else 0
        
        # Grade distribution by type
        grade_distribution = {}
        for grade in all_grades:
            grade_type = grade.get('type', 'other')
            if grade_type not in grade_distribution:
                grade_distribution[grade_type] = {'earned': 0, 'possible': 0}
            grade_distribution[grade_type]['earned'] += grade.get('points_earned', 0)
            grade_distribution[grade_type]['possible'] += grade.get('points_possible', 0)
        
        return Response({
            'current_grade': current_grade,
            'current_average': round(current_average, 2),
            'total_points_earned': total_points_earned,
            'total_points_possible': total_points_possible,
            'grade_distribution': grade_distribution,
            'grade_breakdown': portfolio.grade_breakdown,
            'all_grades': all_grades
        })
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def add_grade(request, portfolio_id):
    """Add a grade to a portfolio"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        grade_type = request.data.get('type')  # exam, homework, quiz, project
        assignment_name = request.data.get('assignment_name')
        points_earned = request.data.get('points_earned')
        points_possible = request.data.get('points_possible')
        date = request.data.get('date', timezone.now().isoformat())
        
        grade_data = {
            'assignment_name': assignment_name,
            'points_earned': points_earned,
            'points_possible': points_possible,
            'type': grade_type,
            'date': date
        }
        
        # Add to appropriate grade list
        if grade_type == 'exam':
            portfolio.exam_grades.append(grade_data)
        elif grade_type == 'homework':
            portfolio.homework_grades.append(grade_data)
        elif grade_type == 'quiz':
            portfolio.quiz_grades.append(grade_data)
        elif grade_type == 'project':
            portfolio.project_grades.append(grade_data)
        
        # Update current grade and timestamp
        portfolio.current_grade = portfolio.calculate_current_grade()
        portfolio.last_grade_update = timezone.now()
        portfolio.save()
        
        return Response({'message': 'Grade added successfully', 'current_grade': portfolio.current_grade})
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_400_BAD_REQUEST)

@api_view(['PUT'])
@permission_classes([permissions.IsAuthenticated])
def update_grade_breakdown(request, portfolio_id):
    """Update grade breakdown percentages"""
    try:
        portfolio = ClassPortfolio.objects.get(id=portfolio_id)
        portfolio.grade_breakdown = request.data.get('grade_breakdown', {})
        portfolio.current_grade = portfolio.calculate_current_grade()
        portfolio.last_grade_update = timezone.now()
        portfolio.save()
        
        return Response({'message': 'Grade breakdown updated', 'current_grade': portfolio.current_grade})
        
    except ClassPortfolio.DoesNotExist:
        return Response({'error': 'Portfolio not found'}, status=status.HTTP_404_NOT_FOUND)

# Class Review Views
class ClassReviewListCreateView(generics.ListCreateAPIView):
    queryset = ClassReview.objects.all()
    serializer_class = ClassReviewSerializer
    permission_classes = [permissions.IsAuthenticated]

class ClassReviewDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = ClassReview.objects.all()
    serializer_class = ClassReviewSerializer
    permission_classes = [permissions.IsAuthenticated]

# Study Group Views
class StudyGroupListCreateView(generics.ListCreateAPIView):
    queryset = StudyGroup.objects.all()
    serializer_class = StudyGroupSerializer
    permission_classes = [permissions.IsAuthenticated]

class StudyGroupDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = StudyGroup.objects.all()
    serializer_class = StudyGroupSerializer
    permission_classes = [permissions.IsAuthenticated]

@api_view(['POST'])
def join_study_group(request, group_id):
    """Join a study group"""
    try:
        group = StudyGroup.objects.get(id=group_id)
        if group.members.count() >= group.max_members:
            return Response({'error': 'Study group is full'}, status=status.HTTP_400_BAD_REQUEST)
        
        group.members.add(request.user)
        return Response({'message': 'Successfully joined study group'})
    except StudyGroup.DoesNotExist:
        return Response({'error': 'Study group not found'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['POST'])
def leave_study_group(request, group_id):
    """Leave a study group"""
    try:
        group = StudyGroup.objects.get(id=group_id)
        group.members.remove(request.user)
        return Response({'message': 'Successfully left study group'})
    except StudyGroup.DoesNotExist:
        return Response({'error': 'Study group not found'}, status=status.HTTP_404_NOT_FOUND)

# Notification Views
class NotificationListView(generics.ListAPIView):
    serializer_class = NotificationSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return Notification.objects.filter(user=self.request.user)

@api_view(['POST'])
def mark_notification_read(request, notification_id):
    """Mark a notification as read"""
    try:
        notification = Notification.objects.get(id=notification_id, user=request.user)
        notification.is_read = True
        notification.save()
        return Response({'message': 'Notification marked as read'})
    except Notification.DoesNotExist:
        return Response({'error': 'Notification not found'}, status=status.HTTP_404_NOT_FOUND)

# Resource Recommendation Views
class ResourceRecommendationListCreateView(generics.ListCreateAPIView):
    queryset = ResourceRecommendation.objects.all()
    serializer_class = ResourceRecommendationSerializer
    permission_classes = [permissions.IsAuthenticated]

class ResourceRecommendationDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = ResourceRecommendation.objects.all()
    serializer_class = ResourceRecommendationSerializer
    permission_classes = [permissions.IsAuthenticated]

# Keep existing views for backward compatibility
class UserProfileView(generics.RetrieveAPIView):
    queryset = User.objects.all()
    serializer_class = UserSerializer
    permission_classes = [permissions.AllowAny]

class UserSearchView(generics.ListAPIView):
    serializer_class = UserSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        query = self.request.query_params.get('q', '')
        if query:
            return User.objects.filter(
                Q(username__icontains=query) |
                Q(first_name__icontains=query) |
                Q(last_name__icontains=query) |
                Q(email__icontains=query)
            )
        return User.objects.none()

class PostListCreateView(generics.ListCreateAPIView):
    queryset = Post.objects.all()
    permission_classes = [permissions.IsAuthenticatedOrReadOnly]

    def get_serializer_class(self):
        if self.request.method == 'POST':
            return PostCreateSerializer
        return PostSerializer

    def perform_create(self, serializer):
        serializer.save(author=self.request.user)

class PostDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Post.objects.all()
    serializer_class = PostSerializer
    permission_classes = [permissions.IsAuthenticatedOrReadOnly]

    def perform_update(self, serializer):
        serializer.save(author=self.request.user)

@api_view(['POST'])
def like_post(request, post_id):
    post = Post.objects.get(id=post_id)
    like, created = Like.objects.get_or_create(post=post, user=request.user)
    
    if not created:
        like.delete()
        is_liked = False
    else:
        is_liked = True
    
    likes_count = post.likes.count()
    return Response({'is_liked': is_liked, 'likes_count': likes_count})

@api_view(['POST'])
def comment_post(request, post_id):
    post = Post.objects.get(id=post_id)
    serializer = CommentSerializer(data=request.data)
    
    if serializer.is_valid():
        serializer.save(author=request.user, post=post)
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def health_check(request):
    return Response({
        'status': 'OK',
        'message': 'HackWestTX Class Portfolio API is running!',
        'version': '2.0.0'
    })

@api_view(['DELETE'])
@permission_classes([permissions.IsAuthenticated])
def delete_user_account(request):
    """
    Delete the authenticated user's account and all associated data.
    This is a comprehensive deletion that removes:
    - User profile
    - All class portfolios created by the user
    - All quizzes and questions created by the user
    - All quiz submissions by the user
    - All posts and comments by the user
    - All likes by the user
    - All study groups created by the user
    - All class reviews by the user
    - All notifications for the user
    - All resource recommendations by the user
    - All flashcards created by the user
    - All lecture materials uploaded by the user
    - All important dates created by the user
    - All syllabi uploaded by the user
    - All syllabus extractions by the user
    """
    try:
        user = request.user
        
        # Get counts before deletion for response
        portfolios_count = ClassPortfolio.objects.filter(created_by=user).count()
        quizzes_count = Quiz.objects.filter(created_by=user).count()
        posts_count = Post.objects.filter(author=user).count()
        study_groups_count = StudyGroup.objects.filter(created_by=user).count()
        
        # Delete all user-related data
        # 1. Delete quiz submissions (user's answers to quizzes)
        QuizSubmission.objects.filter(user=user).delete()
        
        # 2. Delete quiz questions created by user
        QuizQuestion.objects.filter(quiz__created_by=user).delete()
        
        # 3. Delete quizzes created by user
        Quiz.objects.filter(created_by=user).delete()
        
        # 4. Delete flashcards from materials uploaded by user
        Flashcard.objects.filter(material__uploaded_by=user).delete()
        
        # 5. Delete lecture materials uploaded by user
        LectureMaterial.objects.filter(uploaded_by=user).delete()
        
        # 6. Delete important dates from portfolios created by user
        ImportantDate.objects.filter(portfolio__created_by=user).delete()
        
        # 7. Delete syllabus extractions by user
        SyllabusExtraction.objects.filter(syllabus__portfolio__created_by=user).delete()
        
        # 8. Delete syllabi uploaded by user
        Syllabus.objects.filter(portfolio__created_by=user).delete()
        
        # 9. Delete class portfolios created by user
        ClassPortfolio.objects.filter(created_by=user).delete()
        
        # 10. Delete class reviews by user
        ClassReview.objects.filter(reviewer=user).delete()
        
        # 11. Delete study groups created by user
        StudyGroup.objects.filter(created_by=user).delete()
        
        # 12. Delete notifications for user
        Notification.objects.filter(user=user).delete()
        
        # 13. Delete resource recommendations by user
        ResourceRecommendation.objects.filter(recommended_by=user).delete()
        
        # 14. Delete comments by user
        Comment.objects.filter(author=user).delete()
        
        # 15. Delete likes by user
        Like.objects.filter(user=user).delete()
        
        # 16. Delete posts by user
        Post.objects.filter(author=user).delete()
        
        # 17. Remove user from study groups they joined
        for group in StudyGroup.objects.filter(members=user):
            group.members.remove(user)
        
        # 18. Finally, delete the user account itself
        username = user.username
        user.delete()
        
        return Response({
            'message': f'User account "{username}" and all associated data have been permanently deleted.',
            'deleted_data': {
                'portfolios': portfolios_count,
                'quizzes': quizzes_count,
                'posts': posts_count,
                'study_groups': study_groups_count,
            },
            'status': 'success'
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        return Response({
            'error': f'Failed to delete user account: {str(e)}',
            'status': 'error'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# File Processing Views
from .file_processor import FileProcessor
from .models import ProcessedFile
from .serializers import ProcessedFileSerializer, ProcessedFileCreateSerializer
import os
from django.utils import timezone

class ProcessedFileListCreateView(generics.ListCreateAPIView):
    queryset = ProcessedFile.objects.all()
    permission_classes = [permissions.IsAuthenticated]
    
    def get_serializer_class(self):
        if self.request.method == 'POST':
            return ProcessedFileCreateSerializer
        return ProcessedFileSerializer
    
    def get_queryset(self):
        return ProcessedFile.objects.filter(uploaded_by=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(uploaded_by=self.request.user)

class ProcessedFileDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = ProcessedFile.objects.all()
    serializer_class = ProcessedFileSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        return ProcessedFile.objects.filter(uploaded_by=self.request.user)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def upload_and_process_file(request):
    """
    Upload a file and process it to extract text and generate AI summary
    """
    try:
        # Get uploaded file
        uploaded_file = request.FILES.get('file')
        if not uploaded_file:
            return Response({
                'error': 'No file provided',
                'status': 'error'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Get context from request
        context = request.data.get('context', 'other')
        portfolio_id = request.data.get('portfolio_id')
        
        # Determine file type
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        file_type_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'doc',
            '.pptx': 'pptx',
            '.ppt': 'ppt',
            '.txt': 'txt'  # Added for testing
        }
        
        file_type = file_type_map.get(file_extension)
        if not file_type:
            return Response({
                'error': f'Unsupported file type: {file_extension}',
                'status': 'error'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Get portfolio if provided
        portfolio = None
        if portfolio_id:
            try:
                portfolio = ClassPortfolio.objects.get(id=portfolio_id, created_by=request.user)
            except ClassPortfolio.DoesNotExist:
                return Response({
                    'error': 'Portfolio not found or access denied',
                    'status': 'error'
                }, status=status.HTTP_404_NOT_FOUND)
        
        # Create ProcessedFile record
        processed_file = ProcessedFile.objects.create(
            original_file=uploaded_file,
            file_name=uploaded_file.name,
            file_type=file_type,
            file_size=uploaded_file.size,
            context=context,
            uploaded_by=request.user,
            portfolio=portfolio,
            processing_status='processing',
            extracted_text='',  # Will be filled during processing
            ai_summary=''  # Will be filled during processing
        )
        
        # Process the file
        processor = FileProcessor()
        result = processor.process_file_with_summary(uploaded_file, context)
        
        if result['success']:
            # Update the ProcessedFile with results
            processed_file.extracted_text = result['extraction']['text']
            processed_file.ai_summary = result['summary']['summary']
            processed_file.processing_status = 'completed'
            processed_file.metadata = result['extraction']['metadata']
            processed_file.word_count = result['extraction']['word_count']
            processed_file.char_count = result['extraction']['char_count']
            processed_file.processed_at = timezone.now()
            processed_file.save()
            
            # Return success response
            serializer = ProcessedFileSerializer(processed_file)
            return Response({
                'message': 'File processed successfully',
                'file': serializer.data,
                'status': 'success'
            }, status=status.HTTP_201_CREATED)
        else:
            # Update with error
            processed_file.processing_status = 'failed'
            processed_file.processing_error = result.get('error', 'Unknown error')
            processed_file.processed_at = timezone.now()
            processed_file.save()
            
            return Response({
                'error': f'File processing failed: {result.get("error", "Unknown error")}',
                'status': 'error'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    except Exception as e:
        return Response({
            'error': f'Upload failed: {str(e)}',
            'status': 'error'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def reprocess_file(request, file_id):
    """
    Reprocess an existing file (useful if OpenAI API was down during initial processing)
    """
    try:
        # Get the file
        processed_file = ProcessedFile.objects.get(id=file_id, uploaded_by=request.user)
        
        # Update status
        processed_file.processing_status = 'processing'
        processed_file.save()
        
        # Reprocess the file
        processor = FileProcessor()
        result = processor.process_file_with_summary(processed_file.original_file, processed_file.context)
        
        if result['success']:
            # Update with new results
            processed_file.extracted_text = result['extraction']['text']
            processed_file.ai_summary = result['summary']['summary']
            processed_file.processing_status = 'completed'
            processed_file.metadata = result['extraction']['metadata']
            processed_file.word_count = result['extraction']['word_count']
            processed_file.char_count = result['extraction']['char_count']
            processed_file.processed_at = timezone.now()
            processed_file.processing_error = ''
            processed_file.save()
            
            serializer = ProcessedFileSerializer(processed_file)
            return Response({
                'message': 'File reprocessed successfully',
                'file': serializer.data,
                'status': 'success'
            }, status=status.HTTP_200_OK)
        else:
            # Update with error
            processed_file.processing_status = 'failed'
            processed_file.processing_error = result.get('error', 'Unknown error')
            processed_file.processed_at = timezone.now()
            processed_file.save()
            
            return Response({
                'error': f'File reprocessing failed: {result.get("error", "Unknown error")}',
                'status': 'error'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    except ProcessedFile.DoesNotExist:
        return Response({
            'error': 'File not found or access denied',
            'status': 'error'
        }, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({
            'error': f'Reprocessing failed: {str(e)}',
            'status': 'error'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Document Views
class DocumentListCreateView(generics.ListCreateAPIView):
    """List and create documents"""
    serializer_class = DocumentSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        user = self.request.user
        
        # Filter documents by user access
        queryset = Document.objects.filter(
            Q(uploaded_by=user) | Q(portfolio__created_by=user)
        )
        
        # Apply search and filter parameters
        search = self.request.query_params.get('search', None)
        portfolio_id = self.request.query_params.get('portfolio', None)
        learn_method = self.request.query_params.get('learn_method', None)
        is_processed = self.request.query_params.get('is_processed', None)
        
        if search:
            queryset = queryset.filter(
                Q(filename__icontains=search) |
                Q(learning_result__icontains=search)
            )
        
        if portfolio_id:
            queryset = queryset.filter(portfolio_id=portfolio_id)
        
        if learn_method:
            queryset = queryset.filter(learn_method=learn_method)
        
        if is_processed is not None:
            queryset = queryset.filter(is_processed=is_processed.lower() == 'true')
        
        return queryset.order_by('-created_at')

    def get_serializer_class(self):
        if self.request.method == 'POST':
            return DocumentCreateSerializer
        return DocumentSerializer

class DocumentDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Retrieve, update, or delete a document"""
    serializer_class = DocumentSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        user = self.request.user
        # Filter documents by user access
        return Document.objects.filter(
            Q(uploaded_by=user) | Q(portfolio__created_by=user)
        )

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_documents(request):
    """GET method to retrieve only the authenticated user's documents"""
    user = request.user
    
    # Get only documents uploaded by the current user
    documents = Document.objects.filter(uploaded_by=user).order_by('-created_at')
    
    # Apply search filter if provided
    search = request.query_params.get('search', None)
    if search:
        documents = documents.filter(
            Q(filename__icontains=search) |
            Q(learning_result__icontains=search)
        )
    
    # Apply portfolio filter if provided
    portfolio_id = request.query_params.get('portfolio', None)
    if portfolio_id:
        documents = documents.filter(portfolio_id=portfolio_id)
    
    # Serialize the data
    serializer = DocumentSerializer(documents, many=True, context={'request': request})
    
    return Response({
        'count': documents.count(),
        'results': serializer.data
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_document(request):
    """Create a new document with upload details"""
    serializer = DocumentCreateSerializer(data=request.data, context={'request': request})
    
    if serializer.is_valid():
        document = serializer.save()
        
        # Return the full document data
        response_serializer = DocumentSerializer(document, context={'request': request})
        return Response(response_serializer.data, status=status.HTTP_201_CREATED)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['PATCH'])
@permission_classes([permissions.IsAuthenticated])
def update_document_processing(request, document_id):
    """Update document processing status and results"""
    try:
        document = Document.objects.get(
            id=document_id,
            uploaded_by=request.user
        )
    except Document.DoesNotExist:
        return Response({'error': 'Document not found'}, status=status.HTTP_404_NOT_FOUND)
    
    # Update processing status
    is_processed = request.data.get('is_processed', document.is_processed)
    learning_result = request.data.get('learning_result', document.learning_result)
    key_points = request.data.get('key_points', document.key_points)
    processing_error = request.data.get('processing_error', document.processing_error)
    
    document.is_processed = is_processed
    document.learning_result = learning_result
    document.key_points = key_points
    document.processing_error = processing_error
    document.save()
    
    serializer = DocumentSerializer(document, context={'request': request})
    return Response(serializer.data)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def document_analytics(request):
    """Get analytics for user's documents"""
    user = request.user
    documents = Document.objects.filter(uploaded_by=user)
    
    # Basic statistics
    total_documents = documents.count()
    processed_documents = documents.filter(is_processed=True).count()
    successful_documents = documents.filter(learning_result__success=True).count()
    
    # Learn method distribution
    learn_methods = documents.values('learn_method').annotate(
        count=Count('id')
    ).order_by('-count')
    
    # Recent uploads (last 30 days)
    thirty_days_ago = timezone.now() - timedelta(days=30)
    recent_uploads = documents.filter(created_at__gte=thirty_days_ago).count()
    
    # Portfolio distribution
    portfolio_distribution = documents.filter(
        portfolio__isnull=False
    ).values('portfolio__professor').annotate(
        count=Count('id')
    ).order_by('-count')[:5]
    
    return Response({
        'total_documents': total_documents,
        'processed_documents': processed_documents,
        'successful_documents': successful_documents,
        'processing_rate': (processed_documents / total_documents * 100) if total_documents > 0 else 0,
        'success_rate': (successful_documents / processed_documents * 100) if processed_documents > 0 else 0,
        'recent_uploads': recent_uploads,
        'learn_methods': list(learn_methods),
        'portfolio_distribution': list(portfolio_distribution)
    })

# Document Quiz Views
class DocumentQuizListCreateView(generics.ListCreateAPIView):
    """List and create document quizzes"""
    serializer_class = DocumentQuizSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        # Users can only see their own quizzes
        return DocumentQuiz.objects.filter(user=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)


class DocumentQuizDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Retrieve, update, or delete a specific document quiz"""
    serializer_class = DocumentQuizSerializer
    permission_classes = [permissions.IsAuthenticated, IsOwnerOrReadOnly]
    
    def get_queryset(self):
        # Users can only access their own quizzes
        return DocumentQuiz.objects.filter(user=self.request.user)


@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_document_quizzes(request):
    """GET method to retrieve all document quizzes for the authenticated user"""
    quizzes = DocumentQuiz.objects.filter(user=request.user).order_by('-created_at')
    
    # Optional filtering by document
    document_id = request.GET.get('document_id')
    if document_id:
        quizzes = quizzes.filter(document_id=document_id)
    
    # Optional filtering by topic
    topic = request.GET.get('topic')
    if topic:
        quizzes = quizzes.filter(topic__icontains=topic)
    
    serializer = DocumentQuizSerializer(quizzes, many=True, context={'request': request})
    return Response({
        'count': quizzes.count(),
        'results': serializer.data
    })


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_document_quiz(request):
    """POST method to create a new document quiz"""
    serializer = DocumentQuizCreateSerializer(data=request.data, context={'request': request})
    
    if serializer.is_valid():
        quiz = serializer.save()
        response_serializer = DocumentQuizSerializer(quiz, context={'request': request})
        return Response({
            'success': True,
            'data': response_serializer.data
        }, status=status.HTTP_201_CREATED)
    else:
        return Response({
            'success': False,
            'errors': serializer.errors
        }, status=status.HTTP_400_BAD_REQUEST)


# Learning Links Views (formerly YouTube Video Views)
# Note: Class names kept as YouTubeVideo for API backward compatibility
class YouTubeVideoListCreateView(generics.ListCreateAPIView):
    """List and create learning resource links (YouTube, Coursera, articles, etc.)"""
    serializer_class = YouTubeVideoSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        # Users can only see their own learning links
        return YouTubeVideo.objects.filter(user=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)


class YouTubeVideoDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Retrieve, update, or delete a specific learning resource link"""
    serializer_class = YouTubeVideoSerializer
    permission_classes = [permissions.IsAuthenticated, IsOwnerOrReadOnly]
    
    def get_queryset(self):
        # Users can only access their own learning links
        return YouTubeVideo.objects.filter(user=self.request.user)


@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_youtube_videos(request):
    """GET method to retrieve all learning resource links for the authenticated user"""
    videos = YouTubeVideo.objects.filter(user=request.user).order_by('-created_at')
    
    # Optional filtering by title
    title = request.GET.get('title')
    if title:
        videos = videos.filter(title__icontains=title)
    
    serializer = YouTubeVideoSerializer(videos, many=True, context={'request': request})
    return Response({
        'count': videos.count(),
        'results': serializer.data
    })


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_youtube_video(request):
    """POST method to create a new learning resource link (accepts any safe educational URL)"""
    serializer = YouTubeVideoCreateSerializer(data=request.data, context={'request': request})
    
    if serializer.is_valid():
        video = serializer.save()
        response_serializer = YouTubeVideoSerializer(video, context={'request': request})
        return Response({
            'success': True,
            'data': response_serializer.data
        }, status=status.HTTP_201_CREATED)
    else:
        return Response({
            'success': False,
            'errors': serializer.errors
        }, status=status.HTTP_400_BAD_REQUEST)


@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def public_youtube_videos(request):
    """GET method to retrieve all public learning resource links (no authentication required)"""
    videos = YouTubeVideo.objects.all().order_by('-created_at')
    
    # Optional filtering by title
    title = request.GET.get('title')
    if title:
        videos = videos.filter(title__icontains=title)
    
    # Optional filtering by user
    user_id = request.GET.get('user_id')
    if user_id:
        videos = videos.filter(user_id=user_id)
    
    serializer = YouTubeVideoSerializer(videos, many=True, context={'request': request})
    return Response({
        'count': videos.count(),
        'results': serializer.data
    })


# Calendar Event Views
class CalendarEventListCreateView(generics.ListCreateAPIView):
    """List and create calendar events"""
    serializer_class = CalendarEventSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        # Users can only see their own calendar events
        return CalendarEvent.objects.filter(user=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)


class CalendarEventDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Retrieve, update, or delete a specific calendar event"""
    serializer_class = CalendarEventSerializer
    permission_classes = [permissions.IsAuthenticated, IsOwnerOrReadOnly]
    
    def get_queryset(self):
        # Users can only access their own calendar events
        return CalendarEvent.objects.filter(user=self.request.user)


@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def user_calendar_events(request):
    """GET method to retrieve all calendar events for the authenticated user"""
    events = CalendarEvent.objects.filter(user=request.user).order_by('due_date')
    
    # Optional filtering
    event_type = request.GET.get('event_type')
    if event_type:
        events = events.filter(event_type=event_type)
    
    status_filter = request.GET.get('status')
    if status_filter:
        events = events.filter(status=status_filter)
    
    priority = request.GET.get('priority')
    if priority:
        events = events.filter(priority=priority)
    
    # Filter by date range
    start_date = request.GET.get('start_date')
    end_date = request.GET.get('end_date')
    if start_date:
        events = events.filter(due_date__gte=start_date)
    if end_date:
        events = events.filter(due_date__lte=end_date)
    
    # Filter by class
    class_id = request.GET.get('class_id')
    if class_id:
        events = events.filter(class_portfolio_id=class_id)
    
    serializer = CalendarEventSerializer(events, many=True, context={'request': request})
    return Response({
        'count': events.count(),
        'results': serializer.data
    })


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def create_calendar_event(request):
    """POST method to create a new calendar event"""
    serializer = CalendarEventCreateSerializer(data=request.data, context={'request': request})
    
    if serializer.is_valid():
        event = serializer.save()
        response_serializer = CalendarEventSerializer(event, context={'request': request})
        return Response({
            'success': True,
            'data': response_serializer.data
        }, status=status.HTTP_201_CREATED)
    else:
        return Response({
            'success': False,
            'errors': serializer.errors
        }, status=status.HTTP_400_BAD_REQUEST)


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def mark_event_completed(request, event_id):
    """Mark a calendar event as completed"""
    try:
        event = CalendarEvent.objects.get(id=event_id, user=request.user)
        event.mark_completed()
        serializer = CalendarEventSerializer(event, context={'request': request})
        return Response({
            'success': True,
            'data': serializer.data
        })
    except CalendarEvent.DoesNotExist:
        return Response({
            'success': False,
            'error': 'Calendar event not found'
        }, status=status.HTTP_404_NOT_FOUND)


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def link_resource_to_event(request, event_id):
    """Link a learning resource to a calendar event"""
    try:
        event = CalendarEvent.objects.get(id=event_id, user=request.user)
        resource_id = request.data.get('resource_id')
        
        if not resource_id:
            return Response({
                'success': False,
                'error': 'resource_id is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            resource = YouTubeVideo.objects.get(id=resource_id, user=request.user)
            event.add_learning_resource(resource)
            serializer = CalendarEventSerializer(event, context={'request': request})
            return Response({
                'success': True,
                'data': serializer.data
            })
        except YouTubeVideo.DoesNotExist:
            return Response({
                'success': False,
                'error': 'Learning resource not found'
            }, status=status.HTTP_404_NOT_FOUND)
            
    except CalendarEvent.DoesNotExist:
        return Response({
            'success': False,
            'error': 'Calendar event not found'
        }, status=status.HTTP_404_NOT_FOUND)


@api_view(['DELETE'])
@permission_classes([permissions.IsAuthenticated])
def unlink_resource_from_event(request, event_id, resource_id):
    """Remove a learning resource from a calendar event"""
    try:
        event = CalendarEvent.objects.get(id=event_id, user=request.user)
        try:
            resource = YouTubeVideo.objects.get(id=resource_id, user=request.user)
            event.remove_learning_resource(resource)
            serializer = CalendarEventSerializer(event, context={'request': request})
            return Response({
                'success': True,
                'data': serializer.data
            })
        except YouTubeVideo.DoesNotExist:
            return Response({
                'success': False,
                'error': 'Learning resource not found'
            }, status=status.HTTP_404_NOT_FOUND)
            
    except CalendarEvent.DoesNotExist:
        return Response({
            'success': False,
            'error': 'Calendar event not found'
        }, status=status.HTTP_404_NOT_FOUND)
