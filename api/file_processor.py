"""
File Processing Module for HackWestTX Class Portfolio
Handles PDF, PowerPoint, and Word document text extraction
"""

import os
import io
import logging
from typing import Optional, Dict, Any
from django.core.files.uploadedfile import UploadedFile

# File processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation

# OpenAI for text summarization
import openai
from django.conf import settings

logger = logging.getLogger(__name__)

class FileProcessor:
    """Handles text extraction from various file formats"""
    
    def __init__(self):
        # Initialize OpenAI client only if API key is available
        api_key = getattr(settings, 'OPENAI_API_KEY', None)
        if api_key:
            self.openai_client = openai.OpenAI(api_key=api_key)
        else:
            self.openai_client = None
    
    def extract_text_from_file(self, uploaded_file: UploadedFile) -> Dict[str, Any]:
        """
        Extract text from uploaded file based on its type
        
        Args:
            uploaded_file: Django UploadedFile object
            
        Returns:
            Dict containing extracted text and metadata
        """
        file_name = uploaded_file.name
        file_extension = os.path.splitext(file_name)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return self._extract_from_pdf(uploaded_file)
            elif file_extension in ['.docx', '.doc']:
                return self._extract_from_word(uploaded_file)
            elif file_extension in ['.pptx', '.ppt']:
                return self._extract_from_powerpoint(uploaded_file)
            elif file_extension == '.txt':
                return self._extract_from_text(uploaded_file)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file type: {file_extension}',
                    'text': '',
                    'metadata': {}
                }
        except Exception as e:
            logger.error(f"Error processing file {file_name}: {str(e)}")
            return {
                'success': False,
                'error': f'Error processing file: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_pdf(self, uploaded_file: UploadedFile) -> Dict[str, Any]:
        """Extract text from PDF file"""
        try:
            # Reset file pointer
            uploaded_file.seek(0)
            
            # Create PDF reader
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            
            text_content = []
            metadata = {
                'file_type': 'PDF',
                'page_count': len(pdf_reader.pages),
                'title': '',
                'author': '',
                'subject': '',
                'creator': '',
                'producer': '',
                'creation_date': '',
                'modification_date': ''
            }
            
            # Extract metadata if available
            if pdf_reader.metadata:
                metadata.update({
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'producer': pdf_reader.metadata.get('/Producer', ''),
                    'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')),
                    'modification_date': str(pdf_reader.metadata.get('/ModDate', ''))
                })
            
            # Extract text from each page
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                    continue
            
            full_text = '\n\n'.join(text_content)
            
            return {
                'success': True,
                'text': full_text,
                'metadata': metadata,
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            return {
                'success': False,
                'error': f'Error processing PDF: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_word(self, uploaded_file: UploadedFile) -> Dict[str, Any]:
        """Extract text from Word document"""
        try:
            # Reset file pointer
            uploaded_file.seek(0)
            
            # Create document object
            doc = Document(uploaded_file)
            
            text_content = []
            metadata = {
                'file_type': 'Word Document',
                'paragraph_count': len(doc.paragraphs),
                'title': '',
                'author': '',
                'subject': '',
                'keywords': '',
                'comments': '',
                'last_modified_by': '',
                'created': '',
                'modified': ''
            }
            
            # Extract document properties
            if doc.core_properties:
                metadata.update({
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'keywords': doc.core_properties.keywords or '',
                    'comments': doc.core_properties.comments or '',
                    'last_modified_by': doc.core_properties.last_modified_by or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                })
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(' | '.join(row_text))
                if table_text:
                    text_content.append('\n'.join(table_text))
            
            full_text = '\n\n'.join(text_content)
            
            return {
                'success': True,
                'text': full_text,
                'metadata': metadata,
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
            
        except Exception as e:
            logger.error(f"Error processing Word document: {str(e)}")
            return {
                'success': False,
                'error': f'Error processing Word document: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_powerpoint(self, uploaded_file: UploadedFile) -> Dict[str, Any]:
        """Extract text from PowerPoint presentation"""
        try:
            # Reset file pointer
            uploaded_file.seek(0)
            
            # Create presentation object
            prs = Presentation(uploaded_file)
            
            text_content = []
            metadata = {
                'file_type': 'PowerPoint Presentation',
                'slide_count': len(prs.slides),
                'title': '',
                'author': '',
                'subject': '',
                'keywords': '',
                'comments': '',
                'last_modified_by': '',
                'created': '',
                'modified': ''
            }
            
            # Extract presentation properties
            if prs.core_properties:
                metadata.update({
                    'title': prs.core_properties.title or '',
                    'author': prs.core_properties.author or '',
                    'subject': prs.core_properties.subject or '',
                    'keywords': prs.core_properties.keywords or '',
                    'comments': prs.core_properties.comments or '',
                    'last_modified_by': prs.core_properties.last_modified_by or '',
                    'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                    'modified': str(prs.core_properties.modified) if prs.core_properties.modified else ''
                })
            
            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num + 1} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append('\n'.join(slide_text))
            
            full_text = '\n\n'.join(text_content)
            
            return {
                'success': True,
                'text': full_text,
                'metadata': metadata,
                'word_count': len(full_text.split()),
                'char_count': len(full_text)
            }
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint: {str(e)}")
            return {
                'success': False,
                'error': f'Error processing PowerPoint: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_text(self, uploaded_file: UploadedFile) -> Dict[str, Any]:
        """Extract text from plain text file"""
        try:
            # Reset file pointer
            uploaded_file.seek(0)
            
            # Read the text content
            text_content = uploaded_file.read().decode('utf-8')
            
            metadata = {
                'file_type': 'Text Document',
                'encoding': 'utf-8',
                'line_count': len(text_content.splitlines()),
                'title': '',
                'author': '',
                'created': '',
                'modified': ''
            }
            
            return {
                'success': True,
                'text': text_content,
                'metadata': metadata,
                'word_count': len(text_content.split()),
                'char_count': len(text_content)
            }
            
        except Exception as e:
            logger.error(f"Error processing text file: {str(e)}")
            return {
                'success': False,
                'error': f'Error processing text file: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def summarize_text_with_chatgpt(self, text: str, context: str = "academic content") -> Dict[str, Any]:
        """
        Summarize extracted text using ChatGPT
        
        Args:
            text: The text to summarize
            context: Context about the content (e.g., "lecture notes", "syllabus", "assignment")
            
        Returns:
            Dict containing summary and metadata
        """
        if not self.openai_client:
            return {
                'success': True,
                'summary': 'AI summarization not available - OpenAI API key not configured. Text extraction completed successfully.',
                'metadata': {
                    'model_used': 'none',
                    'original_text_length': len(text),
                    'summary_length': 0,
                    'context': context,
                    'note': 'OpenAI API key not configured'
                }
            }
        
        try:
            # Truncate text if too long (ChatGPT has token limits)
            max_chars = 12000  # Leave room for prompt and response
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[Text truncated due to length]"
            
            # Create the prompt
            prompt = f"""
Please analyze and summarize the following {context}. Provide a comprehensive summary that includes:

1. **Main Topics**: Key subjects and themes covered
2. **Important Points**: Critical information and concepts
3. **Key Details**: Specific facts, dates, or data mentioned
4. **Structure**: How the content is organized
5. **Action Items**: Any tasks, assignments, or deadlines mentioned

Content to analyze:
{text}

Please provide a clear, well-structured summary that would be useful for a student studying this material.
"""
            
            # Call OpenAI API
            response = self.openai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are an expert academic assistant that helps students understand and summarize educational content."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=1000,
                temperature=0.3
            )
            
            summary = response.choices[0].message.content
            
            return {
                'success': True,
                'summary': summary,
                'metadata': {
                    'model_used': 'gpt-3.5-turbo',
                    'original_text_length': len(text),
                    'summary_length': len(summary),
                    'context': context
                }
            }
            
        except Exception as e:
            logger.error(f"Error calling OpenAI API: {str(e)}")
            # Check if it's a quota error
            if "quota" in str(e).lower() or "429" in str(e):
                return {
                    'success': True,
                    'summary': 'AI summarization temporarily unavailable due to API quota limits. Text extraction completed successfully.',
                    'metadata': {
                        'model_used': 'none',
                        'original_text_length': len(text),
                        'summary_length': 0,
                        'context': context,
                        'note': 'OpenAI API quota exceeded'
                    }
                }
            else:
                return {
                    'success': False,
                    'error': f'Error calling OpenAI API: {str(e)}',
                    'summary': '',
                    'metadata': {}
                }
    
    def process_file_with_summary(self, uploaded_file: UploadedFile, context: str = "academic content") -> Dict[str, Any]:
        """
        Complete file processing pipeline: extract text and summarize
        
        Args:
            uploaded_file: Django UploadedFile object
            context: Context about the content for better summarization
            
        Returns:
            Dict containing extracted text, summary, and metadata
        """
        # Extract text from file
        extraction_result = self.extract_text_from_file(uploaded_file)
        
        if not extraction_result['success']:
            return extraction_result
        
        # Summarize the extracted text
        summary_result = self.summarize_text_with_chatgpt(extraction_result['text'], context)
        
        # Combine results
        return {
            'success': True,
            'extraction': extraction_result,
            'summary': summary_result,
            'file_name': uploaded_file.name,
            'file_size': uploaded_file.size,
            'processing_complete': True
        }
